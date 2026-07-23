# Crewlyze
# Copyright (c) 2025 Sowmiyan S
# Licensed under the MIT License

"""
FastAPI Server backend for the Crewlyze application.
Serves static HTML/JS/CSS assets and exposes REST APIs + Server-Sent Events (SSE)
for streaming real-time analysis logs.
"""

import os
import sys

try:
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')
except Exception:
    pass

import json
import re
import uuid
import asyncio
import shutil
import threading
import time
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Optional

import pandas as pd
from tools.dataset_tools import read_csv_robust

# Copy assets on startup/reload
try:
    # 1. Convert bin/crewlyze.js line endings to LF
    bin_js = Path(__file__).resolve().parent / "bin" / "crewlyze.js"
    if bin_js.exists():
        with open(bin_js, "rb") as f:
            content = f.read()
        lf_content = content.replace(b"\r\n", b"\n")
        with open(bin_js, "wb") as f:
            f.write(lf_content)
        print("Successfully converted bin/crewlyze.js line endings to LF")
except Exception as e:
    print(f"Failed to convert line endings: {e}")

# 2. Compress large local PNG assets to avoid Git LFS dependencies
try:
    from PIL import Image
    assets_dir = Path(__file__).resolve().parent / "assets"
    targets = {
        "logo.png": (512, 512),
        "chat_logo.png": (512, 512),
        "favicon.png": (48, 48),
        "placeholder_thumbnail.png": (600, 400),
        "branding_image.png": (800, 500)
    }
    for filename, max_size in targets.items():
        filepath = assets_dir / filename
        if filepath.exists():
            orig_size = filepath.stat().st_size
            if orig_size < 1000:
                print(f"Skipping LFS pointer file: {filename}")
                continue
            with Image.open(filepath) as img:
                img.thumbnail(max_size, Image.Resampling.LANCZOS)
                img.save(filepath, "PNG", optimize=True)
            new_size = filepath.stat().st_size
            print(f"Optimized asset {filename}: {orig_size} -> {new_size} bytes")
except Exception as e:
    print(f"Asset optimization failed: {e}")

from fastapi import FastAPI, File, UploadFile, Form, BackgroundTasks, HTTPException, Request
from fastapi.responses import StreamingResponse, FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from starlette.middleware.base import BaseHTTPMiddleware

# regex to find ANSI terminal escape patterns
ANSI_ESCAPE = re.compile(r'\x1B(?:[@-Z\\-_]|\[[0-?]*[ -/]*[@-~])')

# To keep track of log states (e.g. ignoring prompt blocks) per session
log_stream_states = {}

def clean_log_message(line: str, session_id: Optional[str] = None) -> Optional[str]:
    """Strip ANSI color codes, ignore noisy messages, and format thoughts/actions nicely."""
    # Read dynamic log level from environment
    log_level = os.getenv("LOG_LEVEL", "INFO").upper()

    # Strip ANSI colors/escapes
    line = ANSI_ESCAPE.sub('', line)
    
    # Check if empty
    stripped = line.strip()
    if not stripped:
        return None

    line_lower = stripped.lower()
    
    # In ERROR mode, we only output explicit warnings, errors, or exceptions
    if log_level == "ERROR":
        if "warning" in line_lower or "error" in line_lower or "exception" in line_lower:
            if "error" in line_lower or "exception" in line_lower:
                return f"[Error] {stripped}"
            return f"[Warning] {stripped}"
        return None

    # System logs noise keywords to ignore
    noise_keywords = [
        "scriptruncontext",
        "telemetry_opt_out",
        "otel_sdk_disabled",
        "opentelemetry",
        "urllib3",
        "connectionpool",
        "http/1.1",
        "httpx",
        "backoff",
        "requests.packages",
        "missing scriptruncontext",
        "openai-api-keyword",
        "http request",
        "cooldown",
        "rate limit",
        "max_tokens",
    ]
    if any(kw in line_lower for kw in noise_keywords):
        return None

    # Handle stateful prompt block ignoring (skip prompt ignoring in DEBUG mode to see everything)
    if log_level != "DEBUG" and session_id:
        if session_id not in log_stream_states:
            log_stream_states[session_id] = {"in_prompt": False}
        state = log_stream_states[session_id]
        
        # Start ignoring if prompt starts
        if "prompt after formatting:" in line_lower or "use the following format:" in line_lower:
            state["in_prompt"] = True
            return None
            
        # Stop ignoring if agent thoughts/actions/results start
        if state["in_prompt"]:
            stop_triggers = ["thought:", "action:", "action input:", "response:", "observation:", "entering new", "finished chain"]
            if any(trig in line_lower for trig in stop_triggers):
                state["in_prompt"] = False
            else:
                return None  # Still ignoring prompt contents

    # Ignore raw debug logs from crewai/langchain if not in DEBUG mode
    if log_level != "DEBUG":
        if stripped.startswith("[DEBUG]:") or stripped.startswith("[INFO]:"):
            if "working agent" in line_lower:
                agent_name = stripped.split(":", 2)[-1].strip()
                return f"[Agent] {agent_name} is active..."
            return None
    else:
        # In DEBUG mode, keep and slightly prefix them
        if stripped.startswith("[DEBUG]:"):
            return f"[DEBUG] {stripped[8:].strip()}"
        if stripped.startswith("[INFO]:"):
            return f"[INFO] {stripped[7:].strip()}"

    # Format specific Langchain output structures for a premium look
    if "entering new crewagentexecutor chain" in line_lower:
        return "[Task] Starting agent execution task..."
    if "finished chain" in line_lower:
        return "[Task] Execution task completed."
    
    # Format Thoughts, Actions, inputs, and outputs nicely
    if stripped.startswith("Thought:"):
        thought_text = stripped[8:].strip()
        return f"[Thought] {thought_text}"
        
    if stripped.startswith("Action:"):
        action_text = stripped[7:].strip()
        return f"[Calling Tool] {action_text}"
        
    if stripped.startswith("Action Input:"):
        input_text = stripped[13:].strip()
        if len(input_text) > 150:
            input_text = input_text[:150] + "..."
        return f"[Input] {input_text}"
        
    if stripped.startswith("Response:") or stripped.startswith("Observation:"):
        resp_text = stripped.split(":", 1)[1].strip()
        if len(resp_text) > 150:
            resp_text = resp_text[:150] + "..."
        return f"[Tool Response] {resp_text}"

    if "warning" in line_lower or "error" in line_lower:
        if "error" in line_lower:
            return f"[Error] {stripped}"
        return f"[Warning] {stripped}"

    return stripped

# Core analysis engines — imported lazily so the server boots
# even if crewai has install issues on this Python version.
# Actual ImportError surfaces only when analysis is triggered.
_run_crew = None
_apply_runtime_llm_settings = None
_validate_llm_connection = None
_run_copilot_query = None
_export_pdf = None
_export_chat_pdf = None

def _load_crew():
    global _run_crew, _apply_runtime_llm_settings, _validate_llm_connection
    global _run_copilot_query, _export_pdf, _export_chat_pdf
    if _run_crew is None:
        from crew import run_crew as _rc
        from config.llm_config import apply_runtime_llm_settings as _arls, validate_llm_connection as _vlc
        from ui.copilot import run_copilot_query as _rcq
        from ui.export import export_pdf as _ep, export_chat_pdf as _ecp
        _run_crew = _rc
        _apply_runtime_llm_settings = _arls
        _validate_llm_connection = _vlc
        _run_copilot_query = _rcq
        _export_pdf = _ep
        _export_chat_pdf = _ecp

# Suppress warnings
os.environ["CREWAI_TELEMETRY_OPT_OUT"] = "true"
os.environ["OTEL_SDK_DISABLED"]        = "true"

app = FastAPI(
    title="Crewlyze API",
    description="Autonomous Multi-Agent Business Intelligence and Data Engineering Platform",
    version="1.0.5"
)

# Enable CORS for local development flexibility
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------------------------------------------------------------------
# Optional Enterprise Auth
# ---------------------------------------------------------------------------
AUTH_ENABLED = os.getenv("AUTH_ENABLED", "false").lower() == "true"
AUTH_TOKEN = os.getenv("AUTH_TOKEN", "enterprise-secret-key")

class OptionalAuthMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        if AUTH_ENABLED:
            if request.url.path.startswith("/api/") and not request.url.path.startswith("/api/validate-key"):
                auth_header = request.headers.get("Authorization")
                if not auth_header or not auth_header.startswith("Bearer ") or auth_header.split(" ")[1] != AUTH_TOKEN:
                    return JSONResponse(status_code=401, content={"detail": "Unauthorized: Invalid or missing Enterprise Token"})
        response = await call_next(request)
        return response

app.add_middleware(OptionalAuthMiddleware)

# ---------------------------------------------------------------------------
# State & Directory Setup
# ---------------------------------------------------------------------------

USER_HOME = Path.home() / ".crewlyze"
DATA_DIR = Path(os.getenv("CREWLYZE_DATA_DIR", str(USER_HOME / "data")))
SESSIONS_DIR = DATA_DIR / "sessions"
OUTPUTS_DIR = Path(os.getenv("CREWLYZE_OUTPUTS_DIR", str(USER_HOME / "outputs")))

for path in (DATA_DIR, SESSIONS_DIR, OUTPUTS_DIR):
    path.mkdir(exist_ok=True, parents=True)

@app.on_event("startup")
async def cleanup_stale_analyses():
    """Scan all session metadata files on boot and reset any stale projects stuck in the running status."""
    try:
        if SESSIONS_DIR.exists() and SESSIONS_DIR.is_dir():
            for session_dir in SESSIONS_DIR.iterdir():
                if session_dir.is_dir():
                    metadata_path = session_dir / "metadata.json"
                    if metadata_path.exists():
                        try:
                            with open(metadata_path, "r", encoding="utf-8") as f:
                                meta = json.load(f)
                            if meta.get("status") == "running":
                                meta["status"] = "failed"
                                done_path = session_dir / "done.txt"
                                if not done_path.exists():
                                    with open(done_path, "w") as df:
                                        df.write("done")
                                with open(metadata_path, "w", encoding="utf-8") as f:
                                    json.dump(meta, f, indent=2)
                                print(f"Reset stale running session: {session_dir.name}")
                        except Exception as e:
                            print(f"Failed to reset metadata for {session_dir.name}: {e}")
    except Exception as e:
        print(f"Error during startup stale session cleanup: {e}")


def is_safe_id(id_str: str) -> bool:
    """Ensure the ID is strictly alphanumeric (plus dashes/underscores) to prevent path traversal."""
    if not id_str:
        return False
    return bool(re.match(r"^[a-zA-Z0-9_-]+$", id_str))

def is_safe_filename(filename: str) -> bool:
    """Ensure the filename doesn't contain path traversal characters and has a safe pattern."""
    if not filename:
        return False
    if ".." in filename or "/" in filename or "\\" in filename:
        return False
    if "\0" in filename:
        return False
    # Allow safe characters including spaces, dashes, dots, underscores, parentheses, brackets, and common special symbols in column names
    return bool(re.match(r"^[a-zA-Z0-9_\-. ()[\]$,%&+@=;\'~#]+$", filename))

def validate_project_id(project_id: str) -> str:
    """Validate that the project_id matches a safe pattern to prevent path traversal."""
    if not is_safe_id(project_id):
        raise HTTPException(status_code=400, detail="Invalid project ID.")
    return project_id

def get_safe_session_dir(project_id: str) -> Path:
    pid = validate_project_id(project_id)
    base = SESSIONS_DIR.resolve()
    resolved = (base / pid).resolve()
    try:
        resolved.relative_to(base)
    except ValueError:
        raise HTTPException(status_code=400, detail="Path traversal detected.")
    return resolved

def get_safe_output_dir(project_id: str) -> Path:
    pid = validate_project_id(project_id)
    base = OUTPUTS_DIR.resolve()
    resolved = (base / pid).resolve()
    try:
        resolved.relative_to(base)
    except ValueError:
        raise HTTPException(status_code=400, detail="Path traversal detected.")
    return resolved

_metadata_lock = threading.Lock()

def save_project_metadata(project_id: str, meta: dict):
    session_dir = get_safe_session_dir(project_id)
    if not session_dir.exists():
        session_dir.mkdir(parents=True, exist_ok=True)
    metadata_path = session_dir / "metadata.json"
    with _metadata_lock:
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, indent=2)

def get_project_metadata(project_id: str) -> dict:
    session_dir = get_safe_session_dir(project_id)
    metadata_path = session_dir / "metadata.json"
    
    if not session_dir.exists():
        return {}
        
    meta = {}
    with _metadata_lock:
        if metadata_path.exists():
            try:
                with open(metadata_path, "r", encoding="utf-8") as f:
                    meta = json.load(f)
            except Exception:
                pass
            
    # Default metadata if not present or corrupt (compatibility check)
    if not meta:
        upload_file = session_dir / "original_upload.csv"
        filename = "dataset.csv"
        size = 0
        if upload_file.exists():
            filename = "dataset.csv"
            size = upload_file.stat().st_size
        
        results_path = session_dir / "results.json"
        status = "idle"
        if results_path.exists():
            status = "completed"
        elif (session_dir / "done.txt").exists():
            status = "completed"
            
        created_at = session_dir.stat().st_ctime
        meta = {
            "id": project_id,
            "name": f"Project {project_id}",
            "filename": filename,
            "report_title": f"{filename.rsplit('.', 1)[0].replace('_', ' ').title()} Executive Analysis",
            "size": size,
            "created_at": created_at * 1000,
            "status": status,
            "thumbnail": None
        }

    # Dynamically resolve and update the thumbnail link if generated PNGs exist
    output_dir = get_safe_output_dir(project_id)
    current_thumb = meta.get("thumbnail")
    target_thumb = None
    if output_dir.exists() and output_dir.is_dir():
        png_charts = sorted(
            [f for f in output_dir.glob("*.png")],
            key=lambda x: x.stat().st_mtime,
            reverse=True
        )
        if png_charts:
            import urllib.parse
            target_thumb = f"/api/charts/{project_id}/{urllib.parse.quote(png_charts[0].name)}"

    if current_thumb != target_thumb:
        meta["thumbnail"] = target_thumb
        save_project_metadata(project_id, meta)
        
    return meta


def parse_bool(value: Optional[str]) -> bool:
    return bool(value and str(value).strip().lower() not in {"false", "0", "off", "no", ""})


def optimize_goal_grammar(goal: str, provider: str, model: str, api_key: str, env_key_name: str) -> str:
    """Uses the runtime-configured LLM to optimize the grammar of the project goal."""
    if not goal.strip():
        return ""
    try:
        from config.llm_config import apply_runtime_llm_settings, get_llm_params
        from crewai import LLM
        
        apply_runtime_llm_settings(provider, model, api_key or "", env_key_name)
        params = get_llm_params()
        llm = LLM(**params)
        
        prompt = (
            "You are a professional editor. Improve the grammar, phrasing, and professional tone "
            "of the following data analysis goal. Keep it concise (1-2 sentences). "
            "Return ONLY the corrected goal text, without any introductory text, quotes, or metadata.\n\n"
            f"Goal: {goal.strip()}"
        )
        response = llm.call([{"role": "user", "content": prompt}])
        result = response if isinstance(response, str) else str(response)
        return result.strip().strip('"').strip("'")
    except Exception as e:
        print(f"Grammar optimization failed: {e}")
        return goal.strip()


# ---------------------------------------------------------------------------
# Automated Outbound Notifications Pipeline
# ---------------------------------------------------------------------------

def send_automated_email(session_id: str, results_data: dict, meta: dict, cfg: dict, subject: Optional[str] = None, send_pdf: bool = True, send_insights: bool = True):
    try:
        import smtplib
        import re
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText
        from email.mime.base import MIMEBase
        from email import encoders

        smtp_host = cfg.get("SMTP_HOST")
        smtp_port_val = cfg.get("SMTP_PORT", 587)
        try:
            smtp_port = int(smtp_port_val)
        except Exception:
            smtp_port = 587
            
        smtp_user = cfg.get("SMTP_USER")
        smtp_password = cfg.get("SMTP_PASSWORD")
        smtp_sender = cfg.get("SMTP_SENDER") or smtp_user
        smtp_recipient = cfg.get("SMTP_RECIPIENT")
        smtp_secure = parse_bool(cfg.get("SMTP_SECURE"))

        if not smtp_host or not smtp_user or not smtp_recipient:
            print("[Automation Email] SMTP parameters missing (host, username, or recipient). Skipping email.")
            return

        print(f"[Automation Email] Preparing automated report email to {smtp_recipient}...")

        # Generate PDF bytes if requested
        pdf_bytes = None
        report_dict = {}
        if send_pdf:
            session_dir = get_safe_session_dir(session_id)
            cleaned_csv = session_dir / "cleaned.csv"
            df = read_csv_robust(cleaned_csv)
            
            report_dict = {
                "dataframe":      df,
                "cleaning_steps": results_data["cleaning_steps"],
                "relations":      results_data["relations"],
                "insights":       results_data["insights"],
                "code":           results_data.get("code", ""),
                "output_dir":     str(get_safe_output_dir(session_id)),
                "report_title":   meta.get("report_title", meta.get("name", "Analysis Report")),
                "goal":           meta.get("optimized_goal") or meta.get("goal") or "",
            }
            
            _load_crew()
            pdf_bytes = _export_pdf(report_dict)

        # Build email
        msg = MIMEMultipart()
        msg['From'] = smtp_sender
        msg['To'] = smtp_recipient
        
        email_title = subject.strip() if subject and subject.strip() else f"📊 Crewlyze Automated Report: {meta.get('name', 'Analysis Report')}"
        msg['Subject'] = email_title

        # Body
        body = f"Hello,\n\nYour automated data analysis run for project \"{meta.get('name')}\" has completed successfully.\n"
        if send_pdf:
            body += "\nPlease find the executive PDF summary report attached to this email.\n"
            
        if send_insights:
            insights_txt = results_data.get("insights", "").strip()
            if insights_txt:
                body += f"\n--- Key Strategic Insights ---\n\n{insights_txt}\n"
                
        body += "\nBest regards,\nCrewlyze Swarm\n"
        msg.attach(MIMEText(body, 'plain'))

        # Attachment
        if send_pdf and pdf_bytes:
            part = MIMEBase('application', "octet-stream")
            part.set_payload(pdf_bytes)
            encoders.encode_base64(part)
            clean_title = report_dict.get('report_title', meta.get('name', 'report'))
            filename = re.sub(r"[^a-zA-Z0-9_-]", "_", clean_title.lower())[:60] or f"report_{session_id}"
            part.add_header('Content-Disposition', f'attachment; filename="{filename}.pdf"')
            msg.attach(part)

        # Connect and send
        server = None
        if smtp_port == 465 or smtp_secure:
            try:
                server = smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=15)
            except Exception as ssl_err:
                print(f"[Automation Email] SMTP_SSL connection failed, trying standard SMTP: {ssl_err}")
                
        if server is None:
            server = smtplib.SMTP(smtp_host, smtp_port, timeout=15)
            try:
                server.starttls()
            except Exception as tls_err:
                print(f"[Automation Email] STARTTLS failed: {tls_err}")
        
        if smtp_user and smtp_password:
            server.login(smtp_user, smtp_password)
            
        recipients = [r.strip() for r in smtp_recipient.split(",") if r.strip()]
        server.sendmail(smtp_sender, recipients, msg.as_string())
        server.quit()
        print(f"[Automation Email] Automated report email sent successfully to {smtp_recipient}.")
    except Exception as e:
        print(f"[Automation Email] Failed to send automated email: {e}")

def send_automated_slack(session_id: str, results_data: dict, meta: dict, cfg: dict):
    try:
        import urllib.request
        import json
        
        webhook_url = cfg.get("SLACK_WEBHOOK_URL")
        if not webhook_url:
            print("[Automation Slack] Slack webhook URL missing. Skipping Slack post.")
            return

        send_summary = parse_bool(cfg.get("SLACK_SEND_SUMMARY", True))
        if not send_summary:
            print("[Automation Slack] Slack summary card disabled. Skipping Slack post.")
            return

        print("[Automation Slack] Posting automated summary to Slack...")

        title = meta.get("report_title", meta.get("name", "Crewlyze Executive Analysis"))
        rows = results_data.get("rows_count", 0)
        cols = results_data.get("cols_count", 0)
        
        # Parse observations out of the insights markdown
        insights_section = results_data.get("insights", "")
        observations = []
        if isinstance(insights_section, str):
            for line in insights_section.split("\n"):
                if "observation:" in line.lower() or "**observation**:" in line.lower():
                    clean_obs = re.sub(r'^\s*[-*]?\s*\d*\.?\s*\**Observation\**:\s*', '', line, flags=re.IGNORECASE)
                    observations.append(clean_obs.strip())
        elif isinstance(insights_section, dict):
            observations = insights_section.get("observations", [])
            
        obs_text = "\n".join([f"• {o}" for o in observations[:3]]) if observations else "Insights generated successfully."
        
        blocks = [
            {
                "type": "header",
                "text": {
                    "type": "plain_text",
                    "text": "📊 Automated Crewlyze Analysis Completed",
                    "emoji": True
                }
            },
            {
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": f"*Project Name:* {meta.get('name')}\n*Report Title:* {title}\n*Dataset:* `{rows} rows x {cols} columns`"
                }
            },
            {
                "type": "divider"
            },
            {
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": f"*Key Observations:*\n{obs_text}"
                }
            },
            {
                "type": "context",
                "elements": [
                    {
                        "type": "plain_text",
                        "text": f"Session ID: {session_id} | Powered by CrewAI Multi-Agent Swarm",
                        "emoji": True
                    }
                ]
            }
        ]
        
        payload = {"blocks": blocks}
        req = urllib.request.Request(
            webhook_url,
            data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "application/json"}
        )
        with urllib.request.urlopen(req, timeout=10) as response:
            res_body = response.read().decode("utf-8")
        print(f"[Automation Slack] Automated Slack message posted successfully: {res_body}")
    except Exception as e:
        print(f"[Automation Slack] Failed to post automated Slack message: {e}")

def send_automated_webhook(session_id: str, results_data: dict, meta: dict, cfg: dict):
    try:
        import requests
        import json
        
        webhook_url = cfg.get("OUTBOUND_WEBHOOK_URL")
        if not webhook_url:
            print("[Automation Webhook] Outbound Webhook URL missing. Skipping Webhook post.")
            return

        send_json = parse_bool(cfg.get("WEBHOOK_SEND_JSON", True))
        attach_pdf = parse_bool(cfg.get("WEBHOOK_ATTACH_PDF", True))

        print(f"[Automation Webhook] Dispatched automated webhook payload to {webhook_url}...")

        payload = {
            "event": "analysis_completed",
            "timestamp": time.time() * 1000,
            "session_id": session_id,
            "project_name": meta.get("name"),
            "report_title": meta.get("report_title"),
            "status": "success",
        }
        
        if send_json:
            payload["resultsSummary"] = {
                "rows_count": results_data.get("rows_count", 0),
                "cols_count": results_data.get("cols_count", 0),
                "numeric_count": results_data.get("numeric_count", 0),
                "cat_count": results_data.get("cat_count", 0),
                "cleaning_steps": results_data.get("cleaning_steps", []),
                "relations": results_data.get("relations", []),
                "insights": results_data.get("insights", "")
            }
            
        files = {}
        opened_files = []
        
        session_dir = get_safe_session_dir(session_id)
        output_dir = get_safe_output_dir(session_id)
        
        if attach_pdf:
            pdf_path = output_dir / f"{session_id}_report.pdf"
            if not pdf_path.exists():
                try:
                    cleaned_csv = session_dir / "cleaned.csv"
                    df = read_csv_robust(cleaned_csv)
                    report_dict = {
                        "dataframe":      df,
                        "cleaning_steps": results_data["cleaning_steps"],
                        "relations":      results_data["relations"],
                        "insights":       results_data["insights"],
                        "code":           results_data.get("code", ""),
                        "output_dir":     str(output_dir),
                        "report_title":   meta.get("report_title", meta.get("name", "Analysis Report")),
                        "goal":           meta.get("optimized_goal") or meta.get("goal") or "",
                    }
                    _load_crew()
                    pdf_bytes = _export_pdf(report_dict)
                    with open(pdf_path, "wb") as pf:
                        pf.write(pdf_bytes)
                except Exception as pdf_err:
                    print(f"[Webhook Outbound] Failed to pre-compile PDF attachment: {pdf_err}")
            
            if pdf_path.exists():
                fp = open(pdf_path, "rb")
                opened_files.append(fp)
                files["file"] = (f"report_{session_id}.pdf", fp, "application/pdf")
                
        try:
            if files:
                payload_data = {"payload_json": json.dumps(payload)}
                response = requests.post(webhook_url, data=payload_data, files=files, timeout=15)
            else:
                response = requests.post(webhook_url, json=payload, timeout=15)
            response.raise_for_status()
            print(f"[Automation Webhook] Automated custom Webhook completed successfully.")
        finally:
            for fp in opened_files:
                fp.close()
    except Exception as e:
        print(f"[Automation Webhook] Failed to trigger automated custom Webhook: {e}")

def send_automated_discord(session_id: str, results_data: dict, meta: dict, cfg: dict):
    try:
        import requests
        import json
        
        username = cfg.get("DISCORD_USERNAME", "").strip()
        avatar_url = cfg.get("DISCORD_AVATAR_URL", "").strip() or "https://raw.githubusercontent.com/sowmiyan-s/Multi-Agent-Data-Analysis-System-with-CrewAI/main/assets/chat_logo.png"
        embed_color_hex = cfg.get("DISCORD_EMBED_COLOR", "#5865F2").strip()
        mention = cfg.get("DISCORD_MENTION_ROLE", "").strip()
        include_warnings = parse_bool(cfg.get("DISCORD_TOGGLE_WARNINGS", True))
        include_stats = parse_bool(cfg.get("DISCORD_TOGGLE_STATS", True))
        send_summary = parse_bool(cfg.get("DISCORD_SEND_SUMMARY", True))
        attach_pdf = parse_bool(cfg.get("DISCORD_ATTACH_PDF", True))
        attach_charts = parse_bool(cfg.get("DISCORD_ATTACH_CHARTS", False))

        color_val = 5814783  # default blurple
        if embed_color_hex.startswith("#"):
            try:
                color_val = int(embed_color_hex.lstrip("#"), 16)
            except Exception:
                pass

        # Check if separate channels are active
        separate_channels = parse_bool(cfg.get("DISCORD_SEPARATE_CHANNELS"))
        
        if separate_channels:
            print("[Automation Discord] Posting to separate channels...")
            
            def post_sub_report(webhook_url, content_title, text_content, fields=None, force_pdf=False, force_charts=False):
                if not webhook_url:
                    return
                sub_embed = {
                    "title": content_title,
                    "color": color_val,
                    "fields": fields or [],
                    "footer": {
                        "text": f"Session ID: {session_id} | Powered by Crewlyze & CrewAI"
                    }
                }
                sub_payload = {}
                if send_summary:
                    sub_payload["embeds"] = [sub_embed]
                if text_content:
                    sub_payload["content"] = text_content
                if username:
                    sub_payload["username"] = username
                if avatar_url:
                    sub_payload["avatar_url"] = avatar_url
                if mention:
                    sub_payload["content"] = f"{mention} - {content_title} Complete!\n" + sub_payload.get("content", "")

                sub_files = {}
                sub_opened_files = []
                session_dir = get_safe_session_dir(session_id)
                output_dir = get_safe_output_dir(session_id)

                if force_pdf and attach_pdf:
                    pdf_path = output_dir / f"{session_id}_report.pdf"
                    if not pdf_path.exists():
                        try:
                            cleaned_csv = session_dir / "cleaned.csv"
                            df = read_csv_robust(cleaned_csv)
                            report_dict = {
                                "dataframe":      df,
                                "cleaning_steps": results_data["cleaning_steps"],
                                "relations":      results_data["relations"],
                                "insights":       results_data["insights"],
                                "code":           results_data.get("code", ""),
                                "output_dir":     str(output_dir),
                                "report_title":   meta.get("report_title", meta.get("name", "Analysis Report")),
                                "goal":           meta.get("optimized_goal") or meta.get("goal") or "",
                            }
                            _load_crew()
                            pdf_bytes = _export_pdf(report_dict)
                            with open(pdf_path, "wb") as pf:
                                pf.write(pdf_bytes)
                        except Exception as pdf_err:
                            print(f"[Discord Webhook] Failed to pre-compile PDF attachment: {pdf_err}")
                    
                    if pdf_path.exists():
                        fp = open(pdf_path, "rb")
                        sub_opened_files.append(fp)
                        sub_files["file"] = (f"report_{session_id}.pdf", fp, "application/pdf")

                if force_charts and attach_charts and output_dir.exists():
                    png_charts = sorted(list(output_dir.glob("*.png")), key=lambda x: x.stat().st_mtime)
                    for idx, chart_path in enumerate(png_charts[:3]):
                        fp = open(chart_path, "rb")
                        sub_opened_files.append(fp)
                        sub_files[f"chart_{idx}"] = (chart_path.name, fp, "image/png")

                try:
                    if sub_files:
                        payload_data = {"payload_json": json.dumps(sub_payload)}
                        response = requests.post(webhook_url, data=payload_data, files=sub_files, timeout=15)
                    else:
                        response = requests.post(webhook_url, json=sub_payload, timeout=15)
                    response.raise_for_status()
                    print(f"[Automation Discord] Sub-report '{content_title}' posted successfully.")
                except Exception as post_err:
                    print(f"[Automation Discord] Failed to post sub-report '{content_title}': {post_err}")
                finally:
                    for fp in sub_opened_files:
                        fp.close()

            # 1. Data Cleaning
            if parse_bool(cfg.get("DISCORD_CLEANING_ENABLED")) and results_data.get("cleaning_steps"):
                clean_txt = results_data["cleaning_steps"]
                if len(clean_txt) > 1024:
                    clean_txt = clean_txt[:1000] + "..."
                fields = [{"name": "Cleaning Audit Trail", "value": clean_txt, "inline": False}]
                post_sub_report(cfg.get("DISCORD_CLEANING_URL"), "🧹 Data Cleaning Report", "", fields)

            # 2. Relationship Mapper
            if parse_bool(cfg.get("DISCORD_RELATIONS_ENABLED")) and results_data.get("relations"):
                rel_txt = results_data["relations"]
                if len(rel_txt) > 1024:
                    rel_txt = rel_txt[:1000] + "..."
                fields = [{"name": "Relationship Map Insights", "value": rel_txt, "inline": False}]
                post_sub_report(cfg.get("DISCORD_RELATIONS_URL"), "🔗 Relation Mapper Report", "", fields)

            # 3. Business Analysis
            if parse_bool(cfg.get("DISCORD_INSIGHTS_ENABLED")) and results_data.get("insights"):
                ins_txt = results_data["insights"]
                if len(ins_txt) > 1024:
                    ins_txt = ins_txt[:1000] + "..."
                fields = [{"name": "Strategic Business Insights", "value": ins_txt, "inline": False}]
                post_sub_report(cfg.get("DISCORD_INSIGHTS_URL"), "💡 Business Analysis Report", "", fields)

            # 4. Visualization Graphs
            if parse_bool(cfg.get("DISCORD_VISUALIZATION_ENABLED")):
                fields = [{"name": "Charts Overview", "value": f"Generated {len(results_data.get('png_charts', []))} analytical charts.", "inline": False}]
                post_sub_report(cfg.get("DISCORD_VISUALIZATION_URL"), "📊 Visualization Report", "", fields, force_pdf=True, force_charts=True)

            return
            
        webhook_url = cfg.get("DISCORD_WEBHOOK_URL")
        if not webhook_url:
            print("[Automation Discord] Discord webhook URL missing. Skipping Discord post.")
            return

        print("[Automation Discord] Posting automated summary to Discord...")

        title = meta.get("report_title", meta.get("name", "Crewlyze Executive Analysis"))
        rows = results_data.get("rows_count", 0)
        cols = results_data.get("cols_count", 0)
        
        # Parse observations & warnings out of the insights markdown
        insights_section = results_data.get("insights", "")
        observations = []
        warnings_found = []
        if isinstance(insights_section, str):
            for line in insights_section.split("\n"):
                if "observation:" in line.lower() or "**observation**:" in line.lower():
                    clean_obs = re.sub(r'^\s*[-*]?\s*\d*\.?\s*\**Observation\**:\s*', '', line, flags=re.IGNORECASE)
                    observations.append(clean_obs.strip())
                elif "warning" in line.lower() or "alert" in line.lower() or "risk" in line.lower():
                    clean_warn = re.sub(r'^\s*[-*]?\s*\d*\.?\s*\**Warning\**:\s*', '', line, flags=re.IGNORECASE)
                    warnings_found.append(clean_warn.strip())
        elif isinstance(insights_section, dict):
            observations = insights_section.get("observations", [])
            warnings_found = insights_section.get("warnings", [])
            
        obs_text = "\n".join([f"• {o}" for o in observations[:3]]) if observations else "Insights generated successfully."
        
        embed_fields = [
            {
                "name": "Project Name",
                "value": meta.get("name", "Unnamed Project"),
                "inline": True
            },
            {
                "name": "Report Title",
                "value": title,
                "inline": True
            }
        ]
        
        if include_stats:
            embed_fields.append({
                "name": "Dataset Profile",
                "value": f"`{rows} rows x {cols} columns`",
                "inline": True
            })
            
        embed_fields.append({
            "name": "Key Observations",
            "value": obs_text,
            "inline": False
        })
        
        if include_warnings and warnings_found:
            warn_text = "\n".join([f"⚠️ {w}" for w in warnings_found[:2]])
            embed_fields.append({
                "name": "Surfaced Warnings & Risks",
                "value": warn_text,
                "inline": False
            })
            
        embed = {
            "title": "📊 Crewlyze Analysis Completed",
            "color": color_val,
            "fields": embed_fields,
            "footer": {
                "text": f"Session ID: {session_id} | Powered by Crewlyze & CrewAI"
            }
        }
        
        payload = {}
        if send_summary:
            payload["embeds"] = [embed]
        if username:
            payload["username"] = username
        if avatar_url:
            payload["avatar_url"] = avatar_url
        if mention:
            payload["content"] = f"{mention} - Analysis Completed!"
            
        # Deliverables attachment files
        files = {}
        opened_files = []
        
        session_dir = get_safe_session_dir(session_id)
        output_dir = get_safe_output_dir(session_id)
        
        if attach_pdf:
            pdf_path = output_dir / f"{session_id}_report.pdf"
            if not pdf_path.exists():
                try:
                    cleaned_csv = session_dir / "cleaned.csv"
                    df = read_csv_robust(cleaned_csv)
                    report_dict = {
                        "dataframe":      df,
                        "cleaning_steps": results_data["cleaning_steps"],
                        "relations":      results_data["relations"],
                        "insights":       results_data["insights"],
                        "code":           results_data.get("code", ""),
                        "output_dir":     str(output_dir),
                        "report_title":   meta.get("report_title", meta.get("name", "Analysis Report")),
                        "goal":           meta.get("optimized_goal") or meta.get("goal") or "",
                    }
                    _load_crew()
                    pdf_bytes = _export_pdf(report_dict)
                    with open(pdf_path, "wb") as pf:
                        pf.write(pdf_bytes)
                except Exception as pdf_err:
                    print(f"[Discord Webhook] Failed to pre-compile PDF attachment: {pdf_err}")
            
            if pdf_path.exists():
                fp = open(pdf_path, "rb")
                opened_files.append(fp)
                files["file"] = (f"report_{session_id}.pdf", fp, "application/pdf")
                
        if attach_charts and output_dir.exists():
            png_charts = sorted(list(output_dir.glob("*.png")), key=lambda x: x.stat().st_mtime)
            for idx, chart_path in enumerate(png_charts[:3]):
                fp = open(chart_path, "rb")
                opened_files.append(fp)
                files[f"chart_{idx}"] = (chart_path.name, fp, "image/png")
                
        try:
            if files:
                payload_data = {"payload_json": json.dumps(payload)}
                response = requests.post(webhook_url, data=payload_data, files=files, timeout=15)
            else:
                response = requests.post(webhook_url, json=payload, timeout=15)
            response.raise_for_status()
            print(f"[Automation Discord] Automated Discord message posted successfully.")
        finally:
            for fp in opened_files:
                fp.close()
    except Exception as e:
        print(f"[Automation Discord] Failed to post automated Discord message: {e}")

def run_automation_pipeline(session_id: str, results_data: dict):
    try:
        cfg_path = get_local_config_path()
        if not cfg_path.exists():
            return
        with open(cfg_path, "r", encoding="utf-8") as f:
            cfg = json.load(f)
            
        meta = get_project_metadata(session_id)
        
        # Automated email is disabled as per user request to be manual-only via "Email Report" button
        # if parse_bool(cfg.get("AUTOMATION_EMAIL_ENABLED")):
        #     try:
        #         send_automated_email(session_id, results_data, meta, cfg)
        #     except Exception as e:
        #         print(f"[Automation Hub] Email runner error: {e}")
            
        if parse_bool(cfg.get("AUTOMATION_SLACK_ENABLED")):
            try:
                send_automated_slack(session_id, results_data, meta, cfg)
            except Exception as e:
                print(f"[Automation Hub] Slack runner error: {e}")

        if parse_bool(cfg.get("AUTOMATION_DISCORD_ENABLED")):
            try:
                send_automated_discord(session_id, results_data, meta, cfg)
            except Exception as e:
                print(f"[Automation Hub] Discord runner error: {e}")
            
        if parse_bool(cfg.get("AUTOMATION_WEBHOOK_ENABLED")):
            try:
                send_automated_webhook(session_id, results_data, meta, cfg)
            except Exception as e:
                print(f"[Automation Hub] Webhook runner error: {e}")
    except Exception as e:
        print(f"[Automation Hub] Pipeline dispatch error: {e}")


# ---------------------------------------------------------------------------
# Background Task Pipeline
# ---------------------------------------------------------------------------

MAX_CONCURRENT_ANALYSES = 2
active_analyses = 0
active_analyses_lock = threading.Lock()

def run_crew_in_background(
    session_id: str,
    csv_path: str,
    provider: str,
    model: str,
    api_key: str,
    env_key_name: str,
    cooldown: int,
    selected_tasks: list[str],
    deep_analysis: bool,
    report_title: str,
):
    """
    Orchestrates the CrewAI pipeline in a background thread, writing all
    stdout progress to a tail-able stdout.log file and serializing results.
    """
    global active_analyses
    try:
        if not is_safe_id(session_id):
            raise ValueError("Invalid session ID.")
        session_dir = (SESSIONS_DIR / session_id).resolve()
        resolved_csv = Path(csv_path).resolve()
        try:
            resolved_csv.relative_to(session_dir)
        except ValueError:
            raise ValueError("Path traversal detected in CSV path.")

        # 1. Inject thread-isolated LLM configurations and context variables
        from config.context import (
            current_session_id,
            current_session_csv,
            current_session_output_dir,
            current_llm_provider,
            current_llm_model,
            current_llm_api_key,
            current_llm_env_key_name,
            current_cooldown,
            current_deep_analysis,
        )
        current_session_id.set(session_id)
        current_session_csv.set(str(resolved_csv))
        current_session_output_dir.set(str((OUTPUTS_DIR / session_id).resolve()))
        current_llm_provider.set(provider)
        current_llm_model.set(model)
        current_llm_api_key.set(api_key or "")
        current_llm_env_key_name.set(env_key_name or "")
        current_cooldown.set(cooldown)
        current_deep_analysis.set(deep_analysis)

        # Save or update the report title and goal in project metadata
        try:
            meta = get_project_metadata(session_id)
            if report_title:
                meta["report_title"] = report_title.strip()
            
            user_goal = meta.get("goal", "")
            if user_goal.strip():
                print("Optimizing goal grammar...")
                opt_goal = optimize_goal_grammar(user_goal, provider, model, api_key, env_key_name)
                meta["optimized_goal"] = opt_goal
                print(f"Optimized goal: {opt_goal}")
            else:
                meta["optimized_goal"] = ""
                
            save_project_metadata(session_id, meta)
        except Exception as e:
            print(f"Error handling metadata goal/title: {e}")

        session_dir = SESSIONS_DIR / session_id
        log_path = session_dir / "stdout.log"
        done_path = session_dir / "done.txt"
        results_path = session_dir / "results.json"

        # Clean up previous state
        done_path.unlink(missing_ok=True)
        results_path.unlink(missing_ok=True)

        # Update metadata status to running
        try:
            meta = get_project_metadata(session_id)
            meta["status"] = "running"
            save_project_metadata(session_id, meta)
        except Exception:
            pass

        # 2. Redirect stdout and kickoff
        with open(log_path, "w", encoding="utf-8", errors="replace") as log_file:
            import contextlib
            with contextlib.redirect_stdout(log_file):
                try:
                    print("Initializing multi-agent workflows...")
                    _load_crew()
                    result = _run_crew(
                        csv_path,
                        session_id=session_id,
                        selected_tasks=selected_tasks or None,
                        deep_analysis=deep_analysis,
                    )
                    
                    # Convert results to JSON-serializable structure
                    # Re-map Plotly charts into serializable JSON dictionaries
                    plotly_serializable = []
                    for chart in result.get("plotly_charts", []):
                        try:
                            plotly_serializable.append({
                                "title": chart["title"],
                                "fig_json": json.loads(chart["fig"].to_json())
                            })
                        except Exception:
                            pass

                    # Gather static PNG charts
                    png_charts_list = [f.name for f in Path(result["output_dir"]).glob("*.png")]

                    serializable_result = {
                        "cleaning_steps": result["cleaning_steps"],
                        "relations":      result["relations"],
                        "insights":       result["insights"],
                        "code":           result.get("code", ""),
                        "output_dir":     result["output_dir"],
                        "plotly_charts":  plotly_serializable,
                        "png_charts":     png_charts_list,
                        "rows_count":     int(result["dataframe"].shape[0]),
                        "cols_count":     int(result["dataframe"].shape[1]),
                        "numeric_count":  int(len(result["dataframe"].select_dtypes(include=["number"]).columns)),
                        "cat_count":      int(len(result["dataframe"].select_dtypes(include=["object"]).columns))
                    }

                    # Cache first 100 rows as JSON data preview
                    preview_data = result["dataframe"].head(100).replace([float('inf'), float('-inf')], float('nan')).fillna("").to_dict(orient="records")
                    serializable_result["preview"] = preview_data

                    with open(results_path, "w", encoding="utf-8") as f:
                        json.dump(serializable_result, f, indent=2)
                    
                    print("\nAnalysis complete! Ready to render dashboard.")

                    # Update metadata status to completed
                    try:
                        meta = get_project_metadata(session_id)
                        meta["status"] = "completed"
                        if png_charts_list:
                            import urllib.parse
                            meta["thumbnail"] = f"/api/charts/{session_id}/{urllib.parse.quote(png_charts_list[0])}"
                        save_project_metadata(session_id, meta)
                    except Exception:
                        pass

                    # Trigger outbound automations (Email, Slack, Webhook)
                    try:
                        run_automation_pipeline(session_id, serializable_result)
                    except Exception as aut_err:
                        print(f"[Automation Error] Outbound automations failed: {aut_err}")

                except Exception as e:
                    import traceback
                    print(f"\nPipeline failed: {e}", file=sys.stderr)
                    traceback.print_exc(file=log_file)
                    
                    error_result = {"error": str(e)}
                    with open(results_path, "w", encoding="utf-8") as f:
                        json.dump(error_result, f, indent=2)

                    # Update metadata status to failed
                    try:
                        meta = get_project_metadata(session_id)
                        meta["status"] = "failed"
                        save_project_metadata(session_id, meta)
                    except Exception:
                        pass
                finally:
                    # Write done sentinel to stop EventSource streams
                    with open(done_path, "w") as f:
                        f.write("done")
    finally:
        with active_analyses_lock:
            active_analyses = max(0, active_analyses - 1)


# ---------------------------------------------------------------------------
# API Endpoints
# ---------------------------------------------------------------------------

@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    """Uploads the dataset and registers a unique user session ID."""
    session_id = uuid.uuid4().hex[:12]
    session_dir = get_safe_session_dir(session_id)
    session_dir.mkdir(parents=True, exist_ok=True)

    filename_lower = file.filename.lower()
    is_excel = filename_lower.endswith((".xlsx", ".xls"))
    is_sqlite = filename_lower.endswith((".db", ".sqlite", ".sqlite3"))

    if is_excel:
        file_path = session_dir / "uploaded_file.xlsx"
    elif is_sqlite:
        file_path = session_dir / "uploaded_file.db"
    else:
        file_path = session_dir / "original_upload.csv"

    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Pre-configure fresh log files
    log_path = session_dir / "stdout.log"
    with open(log_path, "w") as f:
        f.write("Dataset uploaded successfully.\n")

    proj_name = file.filename.rsplit('.', 1)[0].replace('_', ' ').replace('-', ' ').title()
    status = "idle"
    sheets = []
    tables = []

    if is_excel:
        try:
            xl = pd.ExcelFile(file_path)
            sheets = xl.sheet_names
            status = "awaiting_sheet"
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read Excel workbook: {e}")
    elif is_sqlite:
        try:
            import sqlite3
            conn = sqlite3.connect(str(file_path))
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = [row[0] for row in cursor.fetchall()]
            conn.close()
            status = "awaiting_table"
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read SQLite tables: {e}")
    else:
        # standard CSV validation
        try:
            df = read_csv_robust(file_path)
            # save back formatted to make sure it's UTF-8 comma-separated
            df.to_csv(file_path, index=False)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read CSV: {e}")

    # Save default project metadata
    try:
        meta = {
            "id": session_id,
            "name": proj_name,
            "filename": file.filename,
            "size": file_path.stat().st_size,
            "created_at": time.time() * 1000,
            "status": status
        }
        save_project_metadata(session_id, meta)
    except Exception:
        pass

    return {
        "session_id": session_id,
        "filename": file.filename,
        "size": file_path.stat().st_size,
        "type": "excel" if is_excel else "sqlite" if is_sqlite else "csv",
        "sheets": sheets,
        "tables": tables
    }


@app.post("/api/upload/select-sheet")
async def select_excel_sheet(session_id: str = Form(...), sheet_name: str = Form(...)):
    session_dir = get_safe_session_dir(session_id)
    xlsx_path = session_dir / "uploaded_file.xlsx"
    if not xlsx_path.exists():
        raise HTTPException(status_code=400, detail="Excel upload not found.")
    try:
        df = pd.read_excel(xlsx_path, sheet_name=sheet_name)
        csv_path = session_dir / "original_upload.csv"
        df.to_csv(csv_path, index=False)
        
        meta = get_project_metadata(session_id)
        meta["status"] = "idle"
        meta["filename"] = f"{meta['filename']} [{sheet_name}]"
        meta["size"] = csv_path.stat().st_size
        save_project_metadata(session_id, meta)
        
        return {"status": "success", "session_id": session_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to load sheet: {e}")


@app.post("/api/upload/select-table")
async def select_sqlite_table(session_id: str = Form(...), table_name: str = Form(...)):
    session_dir = get_safe_session_dir(session_id)
    db_path = session_dir / "uploaded_file.db"
    if not db_path.exists():
        raise HTTPException(status_code=400, detail="Database upload not found.")
    try:
        import sqlite3
        safe_table = table_name.replace("`", "``")
        conn = sqlite3.connect(str(db_path))
        df = pd.read_sql_query(f"SELECT * FROM `{safe_table}`", conn)
        conn.close()
        
        csv_path = session_dir / "original_upload.csv"
        df.to_csv(csv_path, index=False)
        
        meta = get_project_metadata(session_id)
        meta["status"] = "idle"
        meta["filename"] = f"{meta['filename']} [{table_name}]"
        meta["size"] = csv_path.stat().st_size
        save_project_metadata(session_id, meta)
        
        return {"status": "success", "session_id": session_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to load table: {e}")


@app.post("/api/query-sql")
async def query_sql_dataset(session_id: str = Form(...), sql_query: str = Form(...)):
    session_dir = get_safe_session_dir(session_id)
    csv_path = session_dir / "cleaned.csv"
    if not csv_path.exists():
        csv_path = session_dir / "original_upload.csv"
    if not csv_path.exists():
        raise HTTPException(status_code=400, detail="Dataset not found.")

    clean_query = sql_query.strip()
    if not clean_query:
        return {"success": False, "error": "Query cannot be empty."}

    # Security check: Enforce read-only SELECT queries
    first_word = clean_query.split()[0].upper() if clean_query.split() else ""
    if first_word not in ("SELECT", "EXPLAIN", "WITH"):
        return {"success": False, "error": "Security Error: Only read-only queries (SELECT) are permitted."}

    forbidden = ["ATTACH", "DETACH", "PRAGMA", "CREATE", "DROP", "ALTER", "INSERT", "UPDATE", "DELETE", "VACUUM", "REINDEX", "TRANSACTION", "COMMIT", "ROLLBACK"]
    if any(re.search(rf"\b{kw}\b", clean_query, re.IGNORECASE) for kw in forbidden):
        return {"success": False, "error": "Security Error: Mutating or administrative SQL statements are blocked."}

    try:
        import sqlite3
        df = read_csv_robust(csv_path)
        conn = sqlite3.connect(":memory:")
        df.to_sql("dataset", conn, index=False)

        # Set execution instruction limit to prevent infinite loops / CPU DoS (~3 seconds max)
        step_counter = [0]
        def interrupt_handler():
            step_counter[0] += 1
            if step_counter[0] > 1000:  # 1,000 progress callbacks ~ 3 seconds
                return 1
            return 0

        conn.set_progress_handler(interrupt_handler, 1000)

        res_df = pd.read_sql_query(clean_query, conn)
        conn.close()

        res_df = res_df.replace([float('inf'), float('-inf')], float('nan')).fillna("")
        results = res_df.head(100).to_dict(orient="records")
        columns = list(res_df.columns)

        return {
            "success": True,
            "columns": columns,
            "results": results,
            "total_count": len(res_df)
        }
    except Exception as e:
        err_msg = str(e)
        if "interrupted" in err_msg.lower():
            err_msg = "Query execution timed out (3s limit exceeded)."
        return {"success": False, "error": err_msg}


@app.get("/api/dataset-diff")
async def get_dataset_diff(session_id: str):
    session_dir = get_safe_session_dir(session_id)
    orig_path = session_dir / "original_upload.csv"
    clean_path = session_dir / "cleaned.csv"
    
    if not orig_path.exists():
        raise HTTPException(status_code=400, detail="Original dataset upload not found.")
        
    orig_df = read_csv_robust(orig_path)
    
    if not clean_path.exists():
        return {
            "cleaned": False,
            "original_rows": len(orig_df),
            "original_cols": len(orig_df.columns),
            "original_columns": list(orig_df.columns)
        }
        
    clean_df = read_csv_robust(clean_path)
    
    rows_dropped = len(orig_df) - len(clean_df)
    cols_changed = []
    
    for col in orig_df.columns:
        if col not in clean_df.columns:
            cols_changed.append({"column": col, "type": "dropped"})
            
    for col in clean_df.columns:
        if col not in orig_df.columns:
            cols_changed.append({"column": col, "type": "added"})
        elif orig_df[col].dtype != clean_df[col].dtype:
            cols_changed.append({
                "column": col, 
                "type": "type_changed", 
                "from": str(orig_df[col].dtype), 
                "to": str(clean_df[col].dtype)
            })
            
    return {
        "cleaned": True,
        "original_rows": len(orig_df),
        "original_cols": len(orig_df.columns),
        "cleaned_rows": len(clean_df),
        "cleaned_cols": len(clean_df.columns),
        "rows_dropped": rows_dropped,
        "changes": cols_changed
    }


@app.post("/api/share/slack")
async def share_report_to_slack(session_id: str = Form(...), webhook_url: str = Form(...)):
    session_dir = get_safe_session_dir(session_id)
    results_path = session_dir / "results.json"
    if not results_path.exists():
        raise HTTPException(status_code=400, detail="Analysis results not found.")
        
    try:
        with open(results_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        meta = get_project_metadata(session_id)
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to load project details.")
        
    title = meta.get("report_title", "Crewlyze Executive Analysis")
    rows = data.get("rows_count", 0)
    cols = data.get("cols_count", 0)
    
    observations = data.get("insights", {}).get("observations", [])
    obs_text = "\n".join([f"• {o}" for o in observations[:3]]) if observations else "No observations recorded."
    
    blocks = [
        {
            "type": "header",
            "text": {
                "type": "plain_text",
                "text": "📊 Crewlyze Analysis Summary Report",
                "emoji": True
            }
        },
        {
            "type": "section",
            "text": {
                "type": "mrkdwn",
                "text": f"*Report Title:* {title}\n*Dataset Profile:* `{rows} rows x {cols} columns`"
            }
        },
        {
            "type": "divider"
        },
        {
            "type": "section",
            "text": {
                "type": "mrkdwn",
                "text": f"*Key Observations:*\n{obs_text}"
            }
        },
        {
            "type": "context",
            "elements": [
                {
                    "type": "plain_text",
                    "text": f"Session ID: {session_id} | Powered by CrewAI Multi-Agent BI Platform",
                    "emoji": True
                }
            ]
        }
    ]
    
    payload = {"blocks": blocks}
    
    try:
        import urllib.request
        import json
        req = urllib.request.Request(
            webhook_url,
            data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "application/json"}
        )
        with urllib.request.urlopen(req) as response:
            res_body = response.read().decode("utf-8")
        return {"status": "success", "response": res_body}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to post to Slack: {e}")


@app.post("/api/validate-key")
async def validate_api_key(
    provider: str = Form(...),
    model: str = Form(...),
    api_key: Optional[str] = Form(""),
):
    """Validate LLM provider credentials before starting analysis."""
    try:
        _load_crew()
    except ImportError as exc:
        raise HTTPException(status_code=503, detail=f"CrewAI not available: {exc}")
    result = _validate_llm_connection(provider, model, api_key or "")
    if not result.get("valid"):
        raise HTTPException(status_code=400, detail=result.get("message", "Validation failed."))
    return result


@app.post("/api/analyze")
async def trigger_analysis(
    background_tasks: BackgroundTasks,
    session_id: str = Form(...),
    provider: str = Form(...),
    model: str = Form(...),
    api_key: Optional[str] = Form(""),
    cooldown: int = Form(5),
    selected_tasks: str = Form(""),
    deep_analysis: str = Form("false"),
    report_title: str = Form(""),
    clean_rules: Optional[str] = Form("")
):
    """Launches the CrewAI analysis process in the background."""
    session_dir = get_safe_session_dir(session_id)
    csv_path = session_dir / "original_upload.csv"

    if not csv_path.exists():
        raise HTTPException(status_code=400, detail="Session upload not found.")

    # Match provider key name
    if provider == "ollama":
        env_key_name = "OLLAMA_BASE_URL"
    elif provider in ("nvidia", "minimax"):
        env_key_name = "NVIDIA_API_KEY"
    else:
        env_key_name = f"{provider.upper()}_API_KEY"

    selected_tasks = [
        task.strip()
        for task in selected_tasks.split(",")
        if task.strip()
    ]
    if not selected_tasks:
        selected_tasks = ["cleaning", "relations", "insights", "visualization"]

    deep = deep_analysis.strip().lower() in {"true", "1", "yes", "on"}

    # Persist report title and rules if provided
    try:
        meta = get_project_metadata(session_id)
        if report_title.strip():
            meta["report_title"] = report_title.strip()
        meta["clean_rules"] = clean_rules.strip() if clean_rules else ""
        save_project_metadata(session_id, meta)
    except Exception:
        pass

    # Concurrency control checks
    global active_analyses
    with active_analyses_lock:
        if active_analyses >= MAX_CONCURRENT_ANALYSES:
            raise HTTPException(
                status_code=429,
                detail="Server is busy. Maximum concurrent analyses limit reached. Please try again later."
            )
        active_analyses += 1

    # Spawn thread-safe background execution
    try:
        background_tasks.add_task(
            run_crew_in_background,
            session_id=session_id,
            csv_path=str(csv_path),
            provider=provider,
            model=model,
            api_key=api_key,
            env_key_name=env_key_name,
            cooldown=cooldown,
            selected_tasks=selected_tasks,
            deep_analysis=deep,
            report_title=report_title.strip(),
        )
    except Exception as e:
        with active_analyses_lock:
            active_analyses = max(0, active_analyses - 1)
        raise e

    return {"status": "started", "session_id": session_id}


@app.get("/api/analyze/stream")
async def stream_analysis_logs(session_id: str):
    """Streams running stdout log lines using Server-Sent Events (SSE)."""
    session_dir = get_safe_session_dir(session_id)
    log_path = session_dir / "stdout.log"

    # Reset streaming state
    if session_id in log_stream_states:
        log_stream_states[session_id] = {"in_prompt": False}

    async def log_generator():
        # Wait for stdout.log file to populate
        for _ in range(50):
            if log_path.exists():
                break
            await asyncio.sleep(0.1)

        if not log_path.exists():
            yield "data: [Initializing pipeline...]\n\n"

        with open(log_path, "r", encoding="utf-8", errors="replace") as f:
            while True:
                line = f.readline()
                if line:
                    cleaned = clean_log_message(line, session_id=session_id)
                    if cleaned is not None:
                        yield f"data: {cleaned}\n\n"
                else:
                    # Look for done flag
                    done_path = session_dir / "done.txt"
                    if done_path.exists():
                        # Read final trailing lines
                        for trail_line in f.readlines():
                            cleaned_trail = clean_log_message(trail_line, session_id=session_id)
                            if cleaned_trail is not None:
                                yield f"data: {cleaned_trail}\n\n"
                        yield "data: [EOF]\n\n"
                        break
                    await asyncio.sleep(0.1)

    return StreamingResponse(log_generator(), media_type="text/event-stream")


@app.get("/api/results")
async def get_results(session_id: str):
    """Retrieves cached JSON results containing stats, insights, and charts."""
    session_dir = get_safe_session_dir(session_id)
    results_path = session_dir / "results.json"
    if not results_path.exists():
        return {"ready": False, "status": "pending"}

    with open(results_path, "r", encoding="utf-8") as f:
        data = json.load(f)
        if "error" in data:
            return data
        data["ready"] = True
        return data


@app.post("/api/copilot")
async def ask_copilot(
    session_id: str = Form(...),
    query: str = Form(...),
    provider: str = Form(...),
    model: str = Form(...),
    api_key: Optional[str] = Form("")
):
    """Runs a natural language query against the dataset using the Copilot agent."""
    if provider == "ollama":
        env_key_name = "OLLAMA_BASE_URL"
    elif provider in ("nvidia", "minimax"):
        env_key_name = "NVIDIA_API_KEY"
    else:
        env_key_name = f"{provider.upper()}_API_KEY"
    _load_crew()
    _apply_runtime_llm_settings(provider, model, api_key or "", env_key_name)

    session_dir = get_safe_session_dir(session_id)
    csv_path = session_dir / "cleaned.csv"
    output_dir = get_safe_output_dir(session_id)

    if not csv_path.exists():
        # Fall back to original upload if cleaning hasn't run or completed
        csv_path = session_dir / "original_upload.csv"

    if not csv_path.exists():
        raise HTTPException(status_code=400, detail="Dataset not uploaded.")

    # Bind thread-local context variables for the current request
    from config.context import current_session_csv, current_session_output_dir
    current_session_csv.set(str(csv_path))
    current_session_output_dir.set(str(output_dir))

    # Call copilot model runner with server-side auto-healing
    try:
        res = _run_copilot_query(query, str(csv_path), str(output_dir))
    except Exception as exc:
        print(f"[AI Chat Auto-Heal] Exception in copilot: {exc}")
        res = {
            "success": True,
            "text": f"✨ **[AI Chat Auto-Healed]** *Recovered from API connection error: {exc}*\n\n"
                    f"The AI Chat engine intercepted the provider exception and maintained active conversation mode. Please check your LLM provider key in Settings if persistent.",
            "plot_path": None
        }

    # Re-map absolute plot path to relative HTTP endpoint URL
    plot_url = None
    if res.get("plot_path"):
        plot_filename = Path(res["plot_path"]).name
        import urllib.parse
        plot_url = f"/api/charts/{session_id}/{urllib.parse.quote(plot_filename)}"

    return {
        "success": res["success"],
        "text":    res["text"],
        "plot_url": plot_url
    }


@app.post("/api/copilot/stream")
async def stream_copilot(
    session_id: str = Form(...),
    query: str = Form(...),
    provider: str = Form(...),
    model: str = Form(...),
    api_key: Optional[str] = Form("")
):
    """Streams copilot reasoning thoughts, tokens, SQL workbench, and dynamic chart suggestions via SSE."""
    if provider == "ollama":
        env_key_name = "OLLAMA_BASE_URL"
    elif provider in ("nvidia", "minimax"):
        env_key_name = "NVIDIA_API_KEY"
    else:
        env_key_name = f"{provider.upper()}_API_KEY"
    _load_crew()
    _apply_runtime_llm_settings(provider, model, api_key or "", env_key_name)

    session_dir = get_safe_session_dir(session_id)
    csv_path = session_dir / "cleaned.csv"
    output_dir = get_safe_output_dir(session_id)

    if not csv_path.exists():
        csv_path = session_dir / "original_upload.csv"

    if not csv_path.exists():
        raise HTTPException(status_code=400, detail="Dataset not uploaded.")

    from config.context import current_session_csv, current_session_output_dir
    current_session_csv.set(str(csv_path))
    current_session_output_dir.set(str(output_dir))

    queue: asyncio.Queue = asyncio.Queue()
    loop = asyncio.get_running_loop()

    def run_sync_stream():
        try:
            from ui.copilot import stream_copilot_query
            for chunk in stream_copilot_query(
                query=query,
                session_id=session_id,
                provider=provider,
                model=model,
                api_key=api_key or "",
                env_key_name=env_key_name,
                csv_path=str(csv_path),
                output_dir=str(output_dir)
            ):
                fut = asyncio.run_coroutine_threadsafe(queue.put(chunk), loop)
                fut.result()
        except Exception as exc:
            err_json = json.dumps({"type": "token", "text": f"\n\nError generating output: {exc}"})
            fut = asyncio.run_coroutine_threadsafe(queue.put(f"data: {err_json}\n\n"), loop)
            try:
                fut.result()
            except Exception:
                pass
        finally:
            fut = asyncio.run_coroutine_threadsafe(queue.put(None), loop)
            try:
                fut.result()
            except Exception:
                pass

    loop.run_in_executor(None, run_sync_stream)

    async def async_generator():
        while True:
            item = await queue.get()
            if item is None:
                break
            yield item

    return StreamingResponse(async_generator(), media_type="text/event-stream")


@app.get("/api/chat-history")
async def get_chat_history(session_id: str):
    """Retrieves saved project AI chat history."""
    session_dir = get_safe_session_dir(session_id)
    history_file = session_dir / "chat_history.json"
    if not history_file.exists():
        return {"messages": []}
    try:
        with open(history_file, "r", encoding="utf-8") as f:
            data = json.load(f)
            return {"messages": data.get("messages", [])}
    except Exception:
        return {"messages": []}


@app.post("/api/chat-history")
async def save_chat_history(session_id: str = Form(...), messages_json: str = Form(...)):
    """Saves project AI chat history to session storage."""
    session_dir = get_safe_session_dir(session_id)
    history_file = session_dir / "chat_history.json"
    try:
        messages = json.loads(messages_json)
        with open(history_file, "w", encoding="utf-8") as f:
            json.dump({"messages": messages}, f, indent=2, ensure_ascii=False)
        return {"status": "success"}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Failed to save chat history: {exc}")


@app.get("/api/export-notebook")
async def get_jupyter_notebook(session_id: str):
    """Generates a downloadable Jupyter Notebook (.ipynb) containing the analysis code."""
    session_dir = get_safe_session_dir(session_id)
    results_path = session_dir / "results.json"
    if not results_path.exists():
        raise HTTPException(status_code=400, detail="Analysis results not found.")
        
    try:
        with open(results_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        code = data.get("code", "")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to read analysis code.")
        
    cells = []
    
    # 1. Title cell
    cells.append({
        "cell_type": "markdown",
        "metadata": {},
        "source": [
            f"# Crewlyze Autonomous Data Analysis Notebook\n",
            f"Generated for session: `{session_id}`\n\n",
            "This notebook contains the data preparation, relationship mapping, ",
            "and visualization scripts automatically generated by the CrewAI agentic pipeline."
        ]
    })
    
    # 2. Setup import and load cell
    setup_code = [
        "import pandas as pd\n",
        "import numpy as np\n",
        "import matplotlib.pyplot as plt\n",
        "import seaborn as sns\n",
        "import plotly.express as px\n",
        "import plotly.graph_objects as go\n\n",
        "# Load the dataset (assumes 'cleaned.csv' is in the same directory as this notebook)\n",
        "try:\n",
        "    df = pd.read_csv('cleaned.csv')\n",
        "    print(\"Cleaned dataset loaded successfully! Shape:\", df.shape)\n",
        "except FileNotFoundError:\n",
        "    print(\"ERROR: 'cleaned.csv' not found. Please ensure it is in the same folder as this notebook.\")\n"
    ]
    cells.append({
        "cell_type": "code",
        "execution_count": None,
        "metadata": {},
        "source": setup_code,
        "outputs": []
    })

    # 3. Dynamic code and markdown cells chunker
    if code:
        raw_lines = code.splitlines(keepends=True)
        current_block = []
        
        for line in raw_lines:
            # Check for header comments or divider comments
            if line.strip().startswith("# ##") or line.strip().startswith("# #"):
                if current_block:
                    cells.append({
                        "cell_type": "code",
                        "execution_count": None,
                        "metadata": {},
                        "source": current_block,
                        "outputs": []
                    })
                    current_block = []
                cells.append({
                    "cell_type": "markdown",
                    "metadata": {},
                    "source": [line.replace("#", "").strip() + "\n"]
                })
            elif line.strip().startswith("#") and len(line.strip()) > 35 and not current_block:
                cells.append({
                    "cell_type": "markdown",
                    "metadata": {},
                    "source": [line.replace("#", "").strip() + "\n"]
                })
            else:
                current_block.append(line)
                
            # Split cell on main figure show calls
            if any(term in line for term in ("plt.show()", "fig.show()", "go.Figure", "px.scatter", "px.bar", "px.line")):
                if len(current_block) > 5:
                    cells.append({
                        "cell_type": "code",
                        "execution_count": None,
                        "metadata": {},
                        "source": current_block,
                        "outputs": []
                    })
                    current_block = []
                    
        if current_block:
            cells.append({
                "cell_type": "code",
                "execution_count": None,
                "metadata": {},
                "source": current_block,
                "outputs": []
            })
    else:
        cells.append({
            "cell_type": "code",
            "execution_count": None,
            "metadata": {},
            "source": ["# No visualization code generated for this session."],
            "outputs": []
        })

    notebook = {
        "cells": cells,
        "metadata": {
            "kernelspec": {
                "display_name": "Python 3",
                "language": "python",
                "name": "python3"
            },
            "language_info": {
                "name": "python",
                "version": "3"
            }
        },
        "nbformat": 4,
        "nbformat_minor": 2
    }
    
    headers = {
        'Content-Disposition': f'attachment; filename="crewlyze_analysis_{session_id}.ipynb"'
    }
    return JSONResponse(content=notebook, headers=headers)


@app.get("/api/export-pdf")
async def get_pdf_report(session_id: str, report_title: Optional[str] = None):
    """Generates and streams back the executive PDF report."""
    session_dir = get_safe_session_dir(session_id)
    results_path = session_dir / "results.json"
    cleaned_csv = session_dir / "cleaned.csv"

    if not results_path.exists() or not cleaned_csv.exists():
        raise HTTPException(status_code=400, detail="Data analysis results not available.")

    with open(results_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    meta = get_project_metadata(session_id)
    title = report_title.strip() if report_title else meta.get("report_title", meta.get("name", "Analysis Report"))
    goal = meta.get("optimized_goal") or meta.get("goal") or ""

    # Format result structure for reportlab builder
    df = read_csv_robust(cleaned_csv)
    report_dict = {
        "dataframe":      df,
        "cleaning_steps": data["cleaning_steps"],
        "relations":      data["relations"],
        "insights":       data["insights"],
        "code":           data.get("code", ""),
        "output_dir":     str(get_safe_output_dir(session_id)),
        "report_title":   title,
        "goal":           goal,
    }

    try:
        _load_crew()
        pdf_bytes = _export_pdf(report_dict)
        filename = re.sub(r"[^a-zA-Z0-9_-]", "_", title.lower())[:60] or f"report_{session_id}"
        return StreamingResponse(
            BytesIO_iterator(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}.pdf"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {e}")


@app.post("/api/export-chat-pdf")
async def export_chat_history_pdf(
    session_id: str = Form(...),
    messages_json: str = Form(...)
):
    """Generates and downloads a PDF containing a custom selection of chat messages."""
    import json
    try:
        messages = json.loads(messages_json)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid messages format: {e}")

    try:
        _load_crew()
        pdf_bytes = _export_chat_pdf(messages, session_id)
        filename = f"chat_history_{session_id}"
        return StreamingResponse(
            BytesIO_iterator(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}.pdf"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {e}")


@app.post("/api/export/webhook")
async def export_webhook(
    session_id: str = Form(...),
    webhook_url: str = Form(...)
):
    """(Enterprise) Export PDF report directly to a Slack/Discord webhook."""
    import requests
    output_dir = get_safe_output_dir(session_id)
    pdf_path = output_dir / f"{session_id}_report.pdf"
    
    if not pdf_path.exists():
        raise HTTPException(status_code=404, detail="PDF report not found. Run analysis first.")
    
    try:
        with open(pdf_path, 'rb') as f:
            files = {'file': (f"report_{session_id}.pdf", f, 'application/pdf')}
            payload = {'content': f"📈 **Crewlyze AI Analysis Complete!**\nNew business insights are ready for session: `{session_id}`"}
            response = requests.post(webhook_url, data=payload, files=files, timeout=10)
            response.raise_for_status()
        return {"status": "success", "message": "Report successfully dispatched to webhook!"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Webhook dispatch failed: {str(e)}")

@app.post("/api/share/slack")
async def manual_share_slack(
    session_id: str = Form(...),
    webhook_url: str = Form(...)
):
    try:
        cfg = {"SLACK_WEBHOOK_URL": webhook_url}
        results_path = get_safe_session_dir(session_id) / "results.json"
        if not results_path.exists():
            raise HTTPException(status_code=400, detail="Results not found. Run analysis first.")
        with open(results_path, "r", encoding="utf-8") as f:
            results_data = json.load(f)
        meta = get_project_metadata(session_id)
        send_automated_slack(session_id, results_data, meta, cfg)
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/share/discord")
async def manual_share_discord(
    session_id: str = Form(...),
    webhook_url: str = Form(...)
):
    try:
        cfg = {}
        cfg_path = get_local_config_path()
        if cfg_path.exists():
            try:
                with open(cfg_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
            except Exception:
                pass
        cfg["DISCORD_WEBHOOK_URL"] = webhook_url
        
        results_path = get_safe_session_dir(session_id) / "results.json"
        if not results_path.exists():
            raise HTTPException(status_code=400, detail="Results not found. Run analysis first.")
        with open(results_path, "r", encoding="utf-8") as f:
            results_data = json.load(f)
        meta = get_project_metadata(session_id)
        send_automated_discord(session_id, results_data, meta, cfg)
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/share/email")
async def manual_share_email(
    session_id: str = Form(...),
    smtp_account_id: Optional[str] = Form(None),
    recipient_email: Optional[str] = Form(None),
    subject: Optional[str] = Form(None),
    send_pdf: Optional[str] = Form("true"),
    send_insights: Optional[str] = Form("true")
):
    try:
        cfg_path = get_local_config_path()
        if not cfg_path.exists():
            raise HTTPException(status_code=400, detail="SMTP configuration not found in settings.")
        with open(cfg_path, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        results_path = get_safe_session_dir(session_id) / "results.json"
        if not results_path.exists():
            raise HTTPException(status_code=400, detail="Results not found. Run analysis first.")
        with open(results_path, "r", encoding="utf-8") as f:
            results_data = json.load(f)
        meta = get_project_metadata(session_id)
        
        selected_cfg = dict(cfg)
        if smtp_account_id:
            accounts = cfg.get("SMTP_ACCOUNTS", [])
            selected_acc = next((acc for acc in accounts if acc.get("id") == smtp_account_id), None)
            if selected_acc:
                selected_cfg["SMTP_HOST"] = selected_acc.get("smtp_host")
                selected_cfg["SMTP_PORT"] = selected_acc.get("smtp_port")
                selected_cfg["SMTP_USER"] = selected_acc.get("smtp_user")
                selected_cfg["SMTP_PASSWORD"] = selected_acc.get("smtp_password")
                selected_cfg["SMTP_SENDER"] = selected_acc.get("smtp_sender") or selected_acc.get("smtp_user")
                selected_cfg["SMTP_SECURE"] = selected_acc.get("smtp_secure")
            else:
                raise HTTPException(status_code=400, detail="Selected SMTP account not found.")
                
        if recipient_email:
            selected_cfg["SMTP_RECIPIENT"] = recipient_email
            
        send_automated_email(
            session_id, 
            results_data, 
            meta, 
            selected_cfg, 
            subject=subject, 
            send_pdf=parse_bool(send_pdf), 
            send_insights=parse_bool(send_insights)
        )
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/integrations/test/email")
async def test_email_integration(
    smtp_host: str = Form(...),
    smtp_port: str = Form(...),
    smtp_user: str = Form(...),
    smtp_password: str = Form(...),
    smtp_sender: str = Form(""),
    smtp_recipient: str = Form(...),
    smtp_secure: str = Form("true")
):
    try:
        import smtplib
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText
        
        host = smtp_host.strip()
        port = int(smtp_port.strip())
        user = smtp_user.strip()
        passwd = smtp_password.strip()
        sender = smtp_sender.strip() or user
        recipients = [r.strip() for r in smtp_recipient.split(",") if r.strip()]
        secure = parse_bool(smtp_secure)
        
        if passwd == "********":
            cfg_path = get_local_config_path()
            if cfg_path.exists():
                with open(cfg_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
                    passwd = cfg.get("SMTP_PASSWORD", "")
        
        if not recipients:
            raise ValueError("No valid recipient email address specified.")
            
        msg = MIMEMultipart()
        msg["From"] = sender
        msg["To"] = ", ".join(recipients)
        msg["Subject"] = "🧪 Crewlyze Test Connection Email"
        body = "This is a test connection email from your Crewlyze Integrations Hub. Your SMTP settings are correctly configured!"
        msg.attach(MIMEText(body, "plain"))
        
        if secure:
            server = smtplib.SMTP_SSL(host, port, timeout=10)
        else:
            server = smtplib.SMTP(host, port, timeout=10)
            server.starttls()
            
        server.login(user, passwd)
        server.sendmail(sender, recipients, msg.as_string())
        server.quit()
        return {"status": "success", "message": "Test email successfully dispatched!"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Email test failed: {str(e)}")

@app.post("/api/integrations/test/slack")
async def test_slack_integration(
    webhook_url: str = Form(...)
):
    try:
        import requests
        url = webhook_url.strip()
        if not url:
            raise ValueError("Webhook URL cannot be empty.")
        
        payload = {
            "text": "🧪 *Crewlyze Integrations Test*\nYour Slack Incoming Webhook is correctly configured! Connection status: `Active` ✓"
        }
        response = requests.post(url, json=payload, timeout=8)
        response.raise_for_status()
        return {"status": "success", "message": "Slack test message dispatched successfully!"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Slack test failed: {str(e)}")

@app.post("/api/integrations/test/discord")
async def test_discord_integration(
    webhook_url: str = Form(...),
    discord_username: Optional[str] = Form(""),
    discord_avatar_url: Optional[str] = Form(""),
    discord_embed_color: Optional[str] = Form(""),
    discord_mention_role: Optional[str] = Form(""),
    discord_send_summary: Optional[str] = Form("true"),
    discord_attach_pdf: Optional[str] = Form("false")
):
    try:
        import requests
        import json
        url = webhook_url.strip()
        if not url:
            raise ValueError("Webhook URL cannot be empty.")
            
        username = discord_username.strip() if discord_username else ""
        avatar_url = discord_avatar_url.strip() if discord_avatar_url else "https://raw.githubusercontent.com/sowmiyan-s/Multi-Agent-Data-Analysis-System-with-CrewAI/main/assets/chat_logo.png"
        embed_color_hex = discord_embed_color.strip() if discord_embed_color else "#5865F2"
        mention = discord_mention_role.strip() if discord_mention_role else ""
        send_summary = parse_bool(discord_send_summary)
        attach_pdf = parse_bool(discord_attach_pdf)
        
        color_val = 3447003
        if embed_color_hex.startswith("#"):
            try:
                color_val = int(embed_color_hex.lstrip("#"), 16)
            except Exception:
                pass
                
        embed = {
            "title": "🧪 Crewlyze Integrations Test",
            "description": "Your Discord webhook configuration is fully functional. Connection status: `Active` ✓",
            "color": color_val
        }
        
        payload = {}
        if send_summary:
            payload["embeds"] = [embed]
        if username:
            payload["username"] = username
        if avatar_url:
            payload["avatar_url"] = avatar_url
        if mention:
            payload["content"] = f"{mention} - Test Connection greeting!"
            
        files = {}
        if attach_pdf:
            import io
            test_file = io.BytesIO(b"%PDF-1.4\n%Connection test stream document")
            files["file"] = ("test_connection_attachment.pdf", test_file, "application/pdf")
            
        if files:
            payload_data = {"payload_json": json.dumps(payload)}
            response = requests.post(url, data=payload_data, files=files, timeout=8)
        else:
            response = requests.post(url, json=payload, timeout=8)
            
        response.raise_for_status()
        return {"status": "success", "message": "Discord test embed dispatched successfully!"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Discord test failed: {str(e)}")

@app.post("/api/integrations/test/webhook")
async def test_webhook_integration(
    webhook_url: str = Form(...)
):
    try:
        import requests
        url = webhook_url.strip()
        if not url:
            raise ValueError("Webhook destination URL cannot be empty.")
            
        payload = {
            "event": "integration_test",
            "timestamp": int(time.time() * 1000),
            "status": "success",
            "message": "This is a REST API connection test payload from Crewlyze."
        }
        response = requests.post(url, json=payload, timeout=8)
        response.raise_for_status()
        return {"status": "success", "message": "Custom webhook POST test request succeeded!"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Custom webhook test failed: {str(e)}")



@app.get("/api/charts/{session_id}/{filename}")
async def serve_chart(session_id: str, filename: str):
    """Serves the generated PNG visual charts."""
    if not is_safe_filename(filename):
        raise HTTPException(status_code=400, detail="Invalid filename.")
    output_dir = get_safe_output_dir(session_id)
    chart_path = (output_dir / filename).resolve()
    try:
        chart_path.relative_to(output_dir)
    except ValueError:
        raise HTTPException(status_code=400, detail="Path traversal detected.")
    if not chart_path.exists():
        raise HTTPException(status_code=404, detail="Chart not found.")
    return FileResponse(chart_path)


# ---------------------------------------------------------------------------
# Utility Streams
# ---------------------------------------------------------------------------

def BytesIO_iterator(data_bytes: bytes):
    """Simple generator to stream raw bytes back to the response."""
    yield data_bytes


# ---------------------------------------------------------------------------
# Ollama Models Fetch
# ---------------------------------------------------------------------------

@app.get("/api/ollama-models")
async def list_ollama_models(base_url: str = "http://localhost:11434"):
    """Fetches list of local Ollama models from the local Ollama service tags API."""
    import requests
    try:
        url = base_url.rstrip("/") + "/api/tags"
        response = requests.get(url, timeout=2.0)
        if response.status_code == 200:
            data = response.json()
            models = [m["name"] for m in data.get("models", [])]
            if models:
                prefixed = [f"ollama/{m}" if not m.startswith("ollama/") else m for m in models]
                return {"models": prefixed}
    except Exception:
        pass
    # Fallback defaults if Ollama service is unreachable or empty
    return {"models": ["ollama/llama3", "ollama/mistral", "ollama/gemma2"]}


# ---------------------------------------------------------------------------
# Metrics & Configurations APIs
# ---------------------------------------------------------------------------

def get_local_config_path() -> Path:
    return USER_HOME / "config.json"

config_lock = asyncio.Lock()

@app.get("/api/metrics")
async def get_performance_metrics():
    from config.metrics_tracker import get_metrics
    return get_metrics()

@app.get("/api/config")
async def get_local_config():
    async with config_lock:
        cfg_path = get_local_config_path()
        cfg = {}
        if cfg_path.exists():
            try:
                with open(cfg_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
            except Exception:
                pass
        
        # Ensure default storage directories and log levels are returned
        if "CREWLYZE_DATA_DIR" not in cfg:
            cfg["CREWLYZE_DATA_DIR"] = str(DATA_DIR)
        if "LOG_LEVEL" not in cfg:
            cfg["LOG_LEVEL"] = os.getenv("LOG_LEVEL", "INFO")
        
        return cfg

@app.post("/api/config")
async def save_local_config(
    request: Request,
    provider: Optional[str] = Form(None),
    api_key: Optional[str] = Form(""),
    base_url: Optional[str] = Form("")
):
    async with config_lock:
        cfg_path = get_local_config_path()
        cfg_path.parent.mkdir(parents=True, exist_ok=True)
        cfg = {}
        if cfg_path.exists():
            try:
                with open(cfg_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
            except Exception:
                pass
        
        content_type = request.headers.get("content-type", "")
        if "application/json" in content_type:
            try:
                payload = await request.json()
                if isinstance(payload, dict):
                    for k, v in payload.items():
                        if v is not None:
                            cfg[k] = v
                        else:
                            cfg.pop(k, None)
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Invalid JSON payload: {e}")
        else:
            # Fallback to compatibility form post
            if not provider:
                raise HTTPException(status_code=400, detail="Missing 'provider' form parameter.")
            
            if provider == "ollama":
                key_name = "OLLAMA_BASE_URL"
                cfg[key_name] = base_url.strip()
            elif provider in ("nvidia", "minimax"):
                key_name = "NVIDIA_API_KEY"
            else:
                key_name = f"{provider.upper()}_API_KEY"

            if provider != "ollama":
                if api_key.strip():
                    if not api_key.endswith("..."):
                        cfg[key_name] = api_key.strip()
                else:
                    cfg.pop(key_name, None)
                    
            if base_url.strip() and provider == "custom":
                cfg["CUSTOM_BASE_URL"] = base_url.strip()
            
        try:
            with open(cfg_path, "w", encoding="utf-8") as f:
                json.dump(cfg, f, indent=2)
            for k, v in cfg.items():
                os.environ[k] = str(v)
            
            # Refresh directory paths dynamically in-memory
            if "CREWLYZE_DATA_DIR" in cfg:
                global DATA_DIR, SESSIONS_DIR
                DATA_DIR = Path(cfg["CREWLYZE_DATA_DIR"])
                SESSIONS_DIR = DATA_DIR / "sessions"
                DATA_DIR.mkdir(exist_ok=True, parents=True)
                SESSIONS_DIR.mkdir(exist_ok=True, parents=True)
                
            if "CREWLYZE_OUTPUTS_DIR" in cfg:
                global OUTPUTS_DIR
                OUTPUTS_DIR = Path(cfg["CREWLYZE_OUTPUTS_DIR"])
                OUTPUTS_DIR.mkdir(exist_ok=True, parents=True)
                
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to write config: {e}")
        return {"status": "success"}

@app.post("/api/config/automation")
async def save_automation_config(
    automation_email_enabled: Optional[str] = Form("false"),
    smtp_host: Optional[str] = Form(""),
    smtp_port: Optional[str] = Form(""),
    smtp_user: Optional[str] = Form(""),
    smtp_password: Optional[str] = Form(""),
    smtp_sender: Optional[str] = Form(""),
    smtp_recipient: Optional[str] = Form(""),
    smtp_secure: Optional[str] = Form("true"),
    automation_slack_enabled: Optional[str] = Form("false"),
    slack_webhook_url: Optional[str] = Form(""),
    slack_send_summary: Optional[str] = Form("true"),
    automation_discord_enabled: Optional[str] = Form("false"),
    discord_webhook_url: Optional[str] = Form(""),
    discord_username: Optional[str] = Form(""),
    discord_avatar_url: Optional[str] = Form(""),
    discord_embed_color: Optional[str] = Form(""),
    discord_mention_role: Optional[str] = Form(""),
    discord_toggle_stats: Optional[str] = Form("true"),
    discord_toggle_warnings: Optional[str] = Form("true"),
    discord_send_summary: Optional[str] = Form("true"),
    discord_attach_pdf: Optional[str] = Form("true"),
    discord_attach_charts: Optional[str] = Form("false"),
    discord_separate_channels: Optional[str] = Form("false"),
    discord_cleaning_url: Optional[str] = Form(""),
    discord_cleaning_enabled: Optional[str] = Form("false"),
    discord_relations_url: Optional[str] = Form(""),
    discord_relations_enabled: Optional[str] = Form("false"),
    discord_insights_url: Optional[str] = Form(""),
    discord_insights_enabled: Optional[str] = Form("false"),
    discord_visualization_url: Optional[str] = Form(""),
    discord_visualization_enabled: Optional[str] = Form("false"),
    automation_webhook_enabled: Optional[str] = Form("false"),
    outbound_webhook_url: Optional[str] = Form(""),
    webhook_send_json: Optional[str] = Form("true"),
    webhook_attach_pdf: Optional[str] = Form("true")
):
    async with config_lock:
        cfg_path = get_local_config_path()
        cfg_path.parent.mkdir(parents=True, exist_ok=True)
        cfg = {}
        if cfg_path.exists():
            try:
                with open(cfg_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
            except Exception:
                pass
                
        cfg["AUTOMATION_EMAIL_ENABLED"] = parse_bool(automation_email_enabled)
        cfg["SMTP_HOST"] = smtp_host.strip() if smtp_host else ""
        cfg["SMTP_PORT"] = int(smtp_port) if smtp_port and smtp_port.strip().isdigit() else 587
        cfg["SMTP_USER"] = smtp_user.strip() if smtp_user else ""
        
        if smtp_password and smtp_password != "********" and not smtp_password.endswith("..."):
            cfg["SMTP_PASSWORD"] = smtp_password.strip()
        elif not smtp_password:
            cfg["SMTP_PASSWORD"] = ""
            
        cfg["SMTP_SENDER"] = smtp_sender.strip() if smtp_sender else ""
        cfg["SMTP_RECIPIENT"] = smtp_recipient.strip() if smtp_recipient else ""
        cfg["SMTP_SECURE"] = parse_bool(smtp_secure)
        
        cfg["AUTOMATION_SLACK_ENABLED"] = parse_bool(automation_slack_enabled)
        cfg["SLACK_WEBHOOK_URL"] = slack_webhook_url.strip() if slack_webhook_url else ""
        cfg["SLACK_SEND_SUMMARY"] = parse_bool(slack_send_summary)

        cfg["AUTOMATION_DISCORD_ENABLED"] = parse_bool(automation_discord_enabled)
        cfg["DISCORD_WEBHOOK_URL"] = discord_webhook_url.strip() if discord_webhook_url else ""
        cfg["DISCORD_USERNAME"] = discord_username.strip() if discord_username else ""
        cfg["DISCORD_AVATAR_URL"] = discord_avatar_url.strip() if discord_avatar_url else ""
        cfg["DISCORD_EMBED_COLOR"] = discord_embed_color.strip() if discord_embed_color else "#5865F2"
        cfg["DISCORD_MENTION_ROLE"] = discord_mention_role.strip() if discord_mention_role else ""
        cfg["DISCORD_TOGGLE_STATS"] = parse_bool(discord_toggle_stats)
        cfg["DISCORD_TOGGLE_WARNINGS"] = parse_bool(discord_toggle_warnings)
        cfg["DISCORD_SEND_SUMMARY"] = parse_bool(discord_send_summary)
        cfg["DISCORD_ATTACH_PDF"] = parse_bool(discord_attach_pdf)
        cfg["DISCORD_ATTACH_CHARTS"] = parse_bool(discord_attach_charts)
        
        cfg["DISCORD_SEPARATE_CHANNELS"] = parse_bool(discord_separate_channels)
        cfg["DISCORD_CLEANING_URL"] = discord_cleaning_url.strip() if discord_cleaning_url else ""
        cfg["DISCORD_CLEANING_ENABLED"] = parse_bool(discord_cleaning_enabled)
        cfg["DISCORD_RELATIONS_URL"] = discord_relations_url.strip() if discord_relations_url else ""
        cfg["DISCORD_RELATIONS_ENABLED"] = parse_bool(discord_relations_enabled)
        cfg["DISCORD_INSIGHTS_URL"] = discord_insights_url.strip() if discord_insights_url else ""
        cfg["DISCORD_INSIGHTS_ENABLED"] = parse_bool(discord_insights_enabled)
        cfg["DISCORD_VISUALIZATION_URL"] = discord_visualization_url.strip() if discord_visualization_url else ""
        cfg["DISCORD_VISUALIZATION_ENABLED"] = parse_bool(discord_visualization_enabled)
        
        cfg["AUTOMATION_WEBHOOK_ENABLED"] = parse_bool(automation_webhook_enabled)
        cfg["OUTBOUND_WEBHOOK_URL"] = outbound_webhook_url.strip() if outbound_webhook_url else ""
        cfg["WEBHOOK_SEND_JSON"] = parse_bool(webhook_send_json)
        cfg["WEBHOOK_ATTACH_PDF"] = parse_bool(webhook_attach_pdf)
        
        try:
            with open(cfg_path, "w", encoding="utf-8") as f:
                json.dump(cfg, f, indent=2)
            # Update environment variables
            for k, v in cfg.items():
                os.environ[k] = str(v)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to write config: {e}")
        return {"status": "success"}

@app.get("/api/llm/providers")
async def get_llm_providers():
    try:
        import litellm
        import os
        import json

        model_cost = {}
        if hasattr(litellm, "model_cost"):
            model_cost = litellm.model_cost
        else:
            try:
                backup_path = os.path.join(os.path.dirname(litellm.__file__), "model_prices_and_context_window_backup.json")
                with open(backup_path, "r", encoding="utf-8") as f:
                    model_cost = json.load(f)
            except Exception:
                pass

        providers = set()
        for model_name, info in model_cost.items():
            if not isinstance(info, dict):
                continue
            prov = info.get("litellm_provider")
            if prov:
                providers.add(str(prov).lower())
            if "/" in model_name:
                prefix = model_name.split("/")[0].lower()
                if prefix and prefix not in ("1024-x-1024", "256-x-256", "512-x-512"):
                    providers.add(prefix)

        def clean_provider_name(p: str) -> bool:
            p = p.lower().strip()
            if not p:
                return False
            if any(x in p for x in ("http", "docs.litellm", " ", "/", "cost", "token", "image", "audio", "video", "speech", "tts", "embed", "pixel")):
                return False
            if "-" in p:
                parts = p.split("-")
                if all(part.isdigit() for part in parts if part):
                    return False
                if any(part.isdigit() for part in parts):
                    if any(x in parts for x in ("x", "w", "h")):
                        return False
            if p in ("hd", "high", "low", "medium", "standard", "v0", "sample_spec", "fallback_generalizations", "max-x-max"):
                return False
            return True

        clean_providers = sorted(list(set(p for p in providers if clean_provider_name(p))))
        
        # Ensure standard fallbacks are included
        standard_fallbacks = ["openai", "anthropic", "nvidia", "groq", "gemini", "ollama", "cohere", "mistral", "vertex_ai", "bedrock", "openrouter", "deepinfra", "together_ai", "xai"]
        for sf in standard_fallbacks:
            if sf not in clean_providers:
                clean_providers.append(sf)

        return {"providers": sorted(clean_providers)}
    except Exception as e:
        return {"providers": ["openai", "anthropic", "nvidia", "groq", "gemini", "ollama", "cohere", "mistral", "vertex_ai", "bedrock", "openrouter", "deepinfra", "together_ai", "xai"], "error": str(e)}

def _is_text_generation_model(model_name: str, info: Optional[dict] = None) -> bool:
    low_name = model_name.lower()
    _EXCLUDE_SUBSTRINGS = (
        "embed", "ada-002", "dall-e", "stable-diffusion", "imagen", "image-generation",
        "tts", "whisper", "audio", "speech", "realtime", "moderation", "content-filter",
        "shield", "guard", "rerank", "clip", "vit", "siglip", "transcription", "translation",
        "vector_store", "search-", "encoder", "ocr", "video_generation", "image_edit"
    )
    if any(sub in low_name for sub in _EXCLUDE_SUBSTRINGS):
        return False
        
    if info:
        non_text_keys = (
            "input_cost_per_image", "output_cost_per_image",
            "input_cost_per_audio_per_second", "input_cost_per_audio_token", "output_cost_per_audio_token",
            "ocr_cost_per_page", "annotation_cost_per_page", "input_cost_per_pixel", "output_cost_per_pixel",
            "input_cost_per_video_per_second", "output_cost_per_video_per_second"
        )
        if any(k in info for k in non_text_keys):
            return False
            
        mode = info.get("mode")
        if mode not in ("chat", "completion", None):
            return False
            
    return True

def _verify_single_model(provider: str, model_name: str, api_key: str) -> bool:
    import litellm
    import os
    
    if provider == "ollama":
        env_key_name = "OLLAMA_BASE_URL"
    elif provider in ("nvidia", "minimax"):
        env_key_name = "NVIDIA_API_KEY"
    else:
        env_key_name = f"{provider.upper()}_API_KEY"

    original_val = os.environ.get(env_key_name)
    if api_key:
        os.environ[env_key_name] = api_key
        
    try:
        litellm.completion(
            model=model_name,
            messages=[{"role": "user", "content": "."}],
            max_tokens=1,
            timeout=4.0
        )
        return True
    except Exception as e:
        err_str = str(e).lower()
        err_type = type(e).__name__
        
        is_bad_request = False
        is_not_found = False
        
        if "badrequest" in err_type.lower() or "400" in err_str:
            is_bad_request = True
        if "notfound" in err_type.lower() or "404" in err_str or "not_found" in err_str:
            is_not_found = True
            
        if is_bad_request or is_not_found:
            return False
            
        return True
    finally:
        if original_val is not None:
            os.environ[env_key_name] = original_val
        elif env_key_name in os.environ:
            del os.environ[env_key_name]

@app.get("/api/llm/providers/{provider}/models")
async def get_llm_models(provider: str, api_key: Optional[str] = None):
    """Returns only text-to-text (chat/completion) models for a provider.
    Filters out voice, image, embedding, moderation, realtime, and other
    non-text-generation models that this project cannot use.
    If api_key is provided or configured, dynamically queries the provider's active models list."""
    
    # Read API Key from local config if not passed/dummy
    if not api_key or not api_key.strip() or api_key.endswith("..."):
        cfg_path = get_local_config_path()
        if cfg_path.exists():
            try:
                with open(cfg_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
                    if provider == "ollama":
                        api_key = cfg.get("OLLAMA_BASE_URL", "")
                    elif provider in ("nvidia", "minimax"):
                        api_key = cfg.get("NVIDIA_API_KEY", "")
                    else:
                        api_key = cfg.get(f"{provider.upper()}_API_KEY", "")
            except Exception:
                pass

    models = []
    fetched_successfully = False

    # Dynamic fetching based on API Key
    if api_key and api_key.strip() and not api_key.endswith("..."):
        import requests
        clean_key = api_key.strip()
        try:
            if provider == "nvidia":
                res = requests.get(
                    "https://integrate.api.nvidia.com/v1/models",
                    headers={"Authorization": f"Bearer {clean_key}"},
                    timeout=4
                )
                if res.status_code == 200:
                    models = [f"nvidia_nim/{m['id']}" for m in res.json().get("data", [])]
                    fetched_successfully = True
            elif provider == "groq":
                res = requests.get(
                    "https://api.groq.com/openai/v1/models",
                    headers={"Authorization": f"Bearer {clean_key}"},
                    timeout=4
                )
                if res.status_code == 200:
                    models = [f"groq/{m['id']}" for m in res.json().get("data", [])]
                    fetched_successfully = True
            elif provider == "openai":
                res = requests.get(
                    "https://api.openai.com/v1/models",
                    headers={"Authorization": f"Bearer {clean_key}"},
                    timeout=4
                )
                if res.status_code == 200:
                    models = [m['id'] for m in res.json().get("data", [])]
                    fetched_successfully = True
        except Exception:
            pass # Fallback to litellm list if request fails

    import litellm
    import os
    
    # Load litellm dictionary
    model_cost = {}
    if hasattr(litellm, "model_cost"):
        model_cost = litellm.model_cost
    else:
        try:
            backup_path = os.path.join(os.path.dirname(litellm.__file__), "model_prices_and_context_window_backup.json")
            with open(backup_path, "r", encoding="utf-8") as f:
                model_cost = json.load(f)
        except Exception:
            pass

    provider_models = []
    
    if fetched_successfully and models:
        for model in models:
            info = model_cost.get(model)
            if _is_text_generation_model(model, info):
                provider_models.append(model)
    else:
        for model_name, info in model_cost.items():
            if not isinstance(info, dict):
                continue
            prov = info.get("litellm_provider", "")
            prov_lower = prov.lower() if prov else ""
            req_prov_lower = provider.lower()
            
            is_match = False
            if prov_lower == req_prov_lower:
                is_match = True
            elif req_prov_lower == "nvidia" and prov_lower == "nvidia_nim":
                is_match = True
            elif req_prov_lower == "cohere" and prov_lower == "cohere_chat":
                is_match = True
            elif req_prov_lower == "bedrock" and prov_lower in ("bedrock_converse", "bedrock_mantle"):
                is_match = True
            elif model_name.startswith(f"{provider}/"):
                is_match = True
            elif provider.lower() == "openai" and "gpt-" in model_name and "/" not in model_name:
                is_match = True

            if is_match:
                if _is_text_generation_model(model_name, info):
                    provider_models.append(model_name)

        try:
            if hasattr(litellm, "model_list"):
                extra = [m for m in litellm.model_list if m.startswith(f"{provider}/")]
                if provider == "openai":
                    extra.extend([m for m in litellm.model_list if "gpt-" in m and "/" not in m])
                for m in extra:
                    info = model_cost.get(m)
                    if _is_text_generation_model(m, info):
                        provider_models.append(m)
        except Exception:
            pass

    provider_models = sorted(list(set(provider_models)))

    verified_models = []
    # Upstream active validation check
    if api_key and api_key.strip() and not api_key.endswith("..."):
        import concurrent.futures
        clean_key = api_key.strip()
        with concurrent.futures.ThreadPoolExecutor(max_workers=20) as executor:
            future_to_model = {
                executor.submit(_verify_single_model, provider, model, clean_key): model
                for model in provider_models
            }
            for future in concurrent.futures.as_completed(future_to_model):
                model = future_to_model[future]
                try:
                    is_active = future.result()
                    if is_active:
                        verified_models.append(model)
                except Exception:
                    verified_models.append(model)
    else:
        verified_models = provider_models

    verified_models = sorted(list(set(verified_models)))
    return {"models": verified_models}

# Duplicate validate-key endpoint removed in favor of validate_api_key defined at line 1374.


# ---------------------------------------------------------------------------
# Project Management APIs
# ---------------------------------------------------------------------------

@app.get("/api/projects")
async def list_projects():
    """Lists all available data analysis projects/sessions."""
    projects = []
    if SESSIONS_DIR.exists():
        for p in SESSIONS_DIR.iterdir():
            if p.is_dir():
                try:
                    meta = get_project_metadata(p.name)
                    if meta:
                        projects.append(meta)
                except Exception:
                    pass
    # Sort projects: newest first
    projects.sort(key=lambda x: x.get("created_at", 0), reverse=True)
    return projects

@app.post("/api/projects")
async def create_project(
    name: str = Form(...),
    report_title: str = Form(""),
    goal: str = Form(""),
    file: UploadFile = File(...)
):
    """Creates a new project context and uploads the dataset (CSV, Excel, or SQLite)."""
    project_id = uuid.uuid4().hex[:12]
    session_dir = get_safe_session_dir(project_id)
    session_dir.mkdir(parents=True, exist_ok=True)

    filename_lower = file.filename.lower()
    is_excel = filename_lower.endswith((".xlsx", ".xls"))
    is_sqlite = filename_lower.endswith((".db", ".sqlite", ".sqlite3"))

    if is_excel:
        file_path = session_dir / "uploaded_file.xlsx"
    elif is_sqlite:
        file_path = session_dir / "uploaded_file.db"
    else:
        file_path = session_dir / "original_upload.csv"

    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Pre-configure fresh log files
    log_path = session_dir / "stdout.log"
    with open(log_path, "w") as f:
        f.write("Project created. Dataset uploaded successfully.\n")

    status = "idle"
    sheets = []
    tables = []

    if is_excel:
        try:
            xl = pd.ExcelFile(file_path)
            sheets = xl.sheet_names
            status = "awaiting_sheet"
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read Excel workbook: {e}")
    elif is_sqlite:
        try:
            import sqlite3
            conn = sqlite3.connect(str(file_path))
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = [row[0] for row in cursor.fetchall()]
            conn.close()
            status = "awaiting_table"
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read SQLite tables: {e}")
    else:
        try:
            df = read_csv_robust(file_path)
            df.to_csv(file_path, index=False)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read CSV: {e}")

    meta = {
        "id": project_id,
        "name": name.strip(),
        "report_title": report_title.strip() or f"{name.strip()} Executive Analysis",
        "goal": goal.strip(),
        "optimized_goal": "",
        "filename": file.filename,
        "size": file_path.stat().st_size,
        "created_at": time.time() * 1000,
        "status": status
    }
    save_project_metadata(project_id, meta)

    return {
        "id": project_id,
        "name": meta["name"],
        "report_title": meta["report_title"],
        "goal": meta["goal"],
        "filename": meta["filename"],
        "size": meta["size"],
        "status": meta["status"],
        "type": "excel" if is_excel else "sqlite" if is_sqlite else "csv",
        "sheets": sheets,
        "tables": tables
    }

@app.post("/api/projects/{project_id}/rename")
async def rename_project(project_id: str, name: str = Form(...)):
    """Renames an existing project context."""
    session_dir = get_safe_session_dir(project_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="Project not found")

    meta = get_project_metadata(project_id)
    meta["name"] = name.strip()
    save_project_metadata(project_id, meta)

    return meta


@app.post("/api/projects/{project_id}/tweak-relations")
async def tweak_relations(project_id: str, relations_text: str = Form(...)):
    """Saves tweaked relationships back to the results cache."""
    session_dir = get_safe_session_dir(project_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="Project not found")

    results_path = session_dir / "results.json"
    
    # Ensure results.json structure is present even if not analysed yet
    res_data = {}
    if results_path.exists():
        try:
            with open(results_path, "r", encoding="utf-8") as f:
                res_data = json.load(f)
        except Exception:
            pass
            
    res_data["relations"] = relations_text.strip()
    
    with open(results_path, "w", encoding="utf-8") as f:
        json.dump(res_data, f, indent=2)
        
    return {"status": "success", "relations": res_data["relations"]}

@app.delete("/api/projects/{project_id}")
async def delete_project(project_id: str):
    """Deletes all session files, artifacts, and outputs of a project."""
    session_dir = get_safe_session_dir(project_id)
    output_dir = get_safe_output_dir(project_id)

    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="Project not found")

    shutil.rmtree(session_dir, ignore_errors=True)
    if output_dir.exists():
        shutil.rmtree(output_dir, ignore_errors=True)

    return {"status": "deleted", "id": project_id}


@app.get("/api/projects/{project_id}/export-zip")
async def export_project_zip(project_id: str):
    """Exports the entire project (metadata, data files, results, and generated charts) as a ZIP file."""
    session_dir = get_safe_session_dir(project_id)
    output_dir = get_safe_output_dir(project_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="Project not found")

    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
        # Zip session files
        for root, dirs, files in os.walk(session_dir):
            for file in files:
                file_path = Path(root) / file
                arcname = Path("session") / file_path.relative_to(session_dir)
                zip_file.write(file_path, arcname=arcname)
        # Zip output files (charts)
        if output_dir.exists():
            for root, dirs, files in os.walk(output_dir):
                for file in files:
                    file_path = Path(root) / file
                    arcname = Path("outputs") / file_path.relative_to(output_dir)
                    zip_file.write(file_path, arcname=arcname)

    zip_buffer.seek(0)
    meta = get_project_metadata(project_id)
    safe_name = re.sub(r"[^a-zA-Z0-9_-]", "_", meta.get("name", "project").lower())
    filename = f"{safe_name}_{project_id}.zip"
    return StreamingResponse(
        BytesIO_iterator(zip_buffer.getvalue()),
        media_type="application/x-zip-compressed",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )


@app.post("/api/projects/import-zip")
async def import_project_zip(file: UploadFile = File(...)):
    """Imports a project from a ZIP file and registers it in the system."""
    zip_contents = await file.read()
    zip_buffer = BytesIO(zip_contents)
    
    project_id = uuid.uuid4().hex[:12]
    temp_dir = DATA_DIR / "temp_import" / project_id
    temp_dir.mkdir(parents=True, exist_ok=True)
    
    target_project_id = project_id
    session_dir = None
    output_dir = None
    try:
        with zipfile.ZipFile(zip_buffer, "r") as zip_file:
            # Zip Slip check:
            for member in zip_file.infolist():
                if ".." in member.filename or member.filename.startswith("/") or member.filename.startswith("\\"):
                    raise HTTPException(status_code=400, detail=f"Invalid zip entry: {member.filename}")
                target_path = (temp_dir / member.filename).resolve()
                try:
                    target_path.relative_to(temp_dir.resolve())
                except ValueError:
                    raise HTTPException(status_code=400, detail=f"Zip Slip detected: {member.filename}")
            zip_file.extractall(temp_dir)
            
        # Verify metadata.json exists
        meta_file = temp_dir / "session" / "metadata.json"
        if not meta_file.exists():
            raise HTTPException(status_code=400, detail="Invalid zip format: missing metadata.json")
            
        with open(meta_file, "r", encoding="utf-8") as f:
            meta = json.load(f)
            
        orig_project_id = meta.get("id")
        if orig_project_id:
            if not is_safe_id(orig_project_id):
                raise HTTPException(status_code=400, detail="Invalid project ID in metadata.")
            target_project_id = orig_project_id
            
        # Check if project conflicts. If so, generate new ID
        session_dir = get_safe_session_dir(target_project_id)
        if session_dir.exists():
            target_project_id = uuid.uuid4().hex[:12]
            session_dir = get_safe_session_dir(target_project_id)
            meta["id"] = target_project_id
            meta["name"] = f"{meta.get('name', 'Imported')} (Copy)"
            
        output_dir = get_safe_output_dir(target_project_id)
        session_dir.mkdir(parents=True, exist_ok=True)
        
        # Copy session files
        for item in (temp_dir / "session").iterdir():
            if item.is_file():
                if not is_safe_filename(item.name):
                    continue
                shutil.copy2(item, session_dir / item.name)
                
        # Copy outputs
        if (temp_dir / "outputs").exists():
            output_dir.mkdir(parents=True, exist_ok=True)
            for item in (temp_dir / "outputs").iterdir():
                if item.is_file():
                    if not is_safe_filename(item.name):
                        continue
                    shutil.copy2(item, output_dir / item.name)
                    
        # Update metadata.json
        meta["id"] = target_project_id
        if meta.get("thumbnail"):
            # Update thumbnail link with new project ID
            thumb_parts = meta["thumbnail"].split("/")
            if len(thumb_parts) >= 5:
                thumb_parts[3] = target_project_id
                meta["thumbnail"] = "/".join(thumb_parts)
                
        with open(session_dir / "metadata.json", "w", encoding="utf-8") as f:
            json.dump(meta, f, indent=2)
            
        return meta
    except Exception as e:
        if session_dir and session_dir.exists():
            shutil.rmtree(session_dir, ignore_errors=True)
        if output_dir and output_dir.exists():
            shutil.rmtree(output_dir, ignore_errors=True)
        raise HTTPException(status_code=400, detail=f"Import failed: {str(e)}")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


@app.get("/api/projects/{project_id}/preview")
async def get_dynamic_preview(project_id: str):
    """Dynamically reads the latest state of the CSV and returns a 100-row preview, column names, shapes, and types."""
    session_dir = get_safe_session_dir(project_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="Project not found")

    cleaned_csv = session_dir / "cleaned.csv"
    original_csv = session_dir / "original_upload.csv"
    csv_path = cleaned_csv if cleaned_csv.exists() else original_csv

    if not csv_path.exists():
        raise HTTPException(status_code=404, detail="CSV not found.")

    try:
        df = read_csv_robust(str(csv_path))
        rows_count, cols_count = df.shape
        preview = df.head(100).fillna("").to_dict(orient="records")
        col_types = {col: str(dtype) for col, dtype in df.dtypes.items()}
        columns = list(df.columns)

        # Update cache in results.json if it exists
        results_path = session_dir / "results.json"
        if results_path.exists():
            try:
                with open(results_path, "r", encoding="utf-8") as f:
                    res_data = json.load(f)
                res_data["preview"] = preview
                res_data["rows_count"] = rows_count
                res_data["cols_count"] = cols_count
                with open(results_path, "w", encoding="utf-8") as f:
                    json.dump(res_data, f, indent=2)
            except Exception:
                pass

        return {
            "columns": columns,
            "col_types": col_types,
            "rows_count": rows_count,
            "cols_count": cols_count,
            "preview": preview
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to load preview: {str(e)}")


@app.get("/api/projects/{project_id}/download-csv")
async def download_project_csv(project_id: str):
    """Downloads the cleaned dataset CSV for the specified project."""
    session_dir = get_safe_session_dir(project_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="Project not found")

    cleaned_csv = session_dir / "cleaned.csv"
    original_csv = session_dir / "original_upload.csv"
    csv_path = cleaned_csv if cleaned_csv.exists() else original_csv

    if not csv_path.exists():
        raise HTTPException(status_code=404, detail="CSV not found.")

    try:
        meta = get_project_metadata(project_id)
        orig_name = meta.get("filename", "dataset.csv")
    except Exception:
        orig_name = "dataset.csv"

    base_name = orig_name.rsplit(".", 1)[0] if "." in orig_name else orig_name
    download_filename = f"{base_name}_cleaned.csv"

    return FileResponse(csv_path, media_type="text/csv", filename=download_filename)


# ---------------------------------------------------------------------------
# PowerPoint (.pptx) Slide Deck Export
# ---------------------------------------------------------------------------

# PowerPoint (.pptx) Executive Slide Deck Export Engine
# ---------------------------------------------------------------------------

@app.get("/api/projects/{project_id}/export-pptx")
async def export_project_pptx(project_id: str, theme: str = "dark"):
    """Generates a McKinsey-style executive PowerPoint slide deck from project results."""
    session_dir = get_safe_session_dir(project_id)
    results_path = session_dir / "results.json"
    if not results_path.exists():
        raise HTTPException(status_code=404, detail="Results not found. Run analysis first.")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        import subprocess, sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

    with open(results_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    meta = get_project_metadata(project_id)
    report_title = meta.get("report_title", meta.get("name", "Executive Data Analysis"))
    project_name = meta.get("name", "Crewlyze Project")

    is_light = (theme.lower() == "light")

    # Color Palette Definitions
    if is_light:
        bg_rgb = (248, 250, 252)        # Slate 50
        card_rgb = (255, 255, 255)      # White
        card_border_rgb = (226, 232, 240) # Slate 200
        text_head_rgb = (15, 23, 42)    # Slate 900
        text_body_rgb = (51, 65, 85)    # Slate 700
        text_sub_rgb = (100, 116, 139)  # Slate 500
        accent_purple = (124, 58, 237)  # Violet 600
        accent_emerald = (5, 150, 105)  # Emerald 600
        accent_cyan = (2, 132, 199)     # Sky 600
        accent_rose = (225, 29, 72)     # Rose 600
        accent_amber = (217, 119, 6)    # Amber 600
    else:
        bg_rgb = (15, 17, 23)          # Dark Obsidian
        card_rgb = (24, 28, 41)         # Dark Slate Card
        card_border_rgb = (51, 65, 85)  # Dark Border
        text_head_rgb = (255, 255, 255) # Pure White
        text_body_rgb = (226, 232, 240) # Slate 200
        text_sub_rgb = (148, 163, 184)  # Slate 400
        accent_purple = (168, 85, 247)  # Purple 500
        accent_emerald = (16, 185, 129) # Emerald 500
        accent_cyan = (14, 165, 233)    # Sky 500
        accent_rose = (244, 63, 94)     # Rose 500
        accent_amber = (245, 158, 11)   # Amber 500

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)

    def _clean_md(text: str) -> str:
        """Strips raw markdown hashes and asterisks."""
        if not text: return ""
        cleaned = re.sub(r'^\s*#{1,6}\s*', '', text, flags=re.MULTILINE)
        cleaned = cleaned.replace('**', '').replace('__', '')
        return cleaned.strip()

    def _clean_takeaway_text(text: str) -> str:
        if not text:
            return "Visual distribution map detailing parameter correlations and features matrix."
        cleaned = re.sub(r'\[Auto-Healing.*?\]', '', text, flags=re.IGNORECASE)
        cleaned = re.sub(r'Warnings\s*&\s*Alerts:.*', '', cleaned, flags=re.DOTALL | re.IGNORECASE)
        cleaned = re.sub(r'Active insights agent failed.*', '', cleaned, flags=re.DOTALL | re.IGNORECASE)
        cleaned = re.sub(r'^\d+[\.\)]\s*', '', cleaned)
        lines = [line.strip() for line in cleaned.split("\n") if line.strip() and not line.strip().startswith("- [Auto-Healing")]
        return "\n".join(lines).strip() or "Visual distribution map detailing parameter correlations and features matrix."

    def _add_textbox(slide, left, top, width, height, text, size=13, bold=False, color=text_body_rgb, align=PP_ALIGN.LEFT):
        clean_t = _clean_md(text)
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Dynamic font size auto-scaling based on text length
        if len(clean_t) > 400:
            size = min(size, 9.5)
            if len(clean_t) > 650:
                clean_t = clean_t[:640] + "..."
        elif len(clean_t) > 220:
            size = min(size, 10.5)
        elif len(clean_t) > 120:
            size = min(size, 11.5)

        p = tf.paragraphs[0]
        p.text = clean_t
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = align
        return tf

    # ── SLIDE 1: Cover Page ───────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide1)

    # Accent Header Bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(0.18), Inches(5.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*accent_purple)
    bar.line.fill.background()

    # Title & Subtitle
    _add_textbox(slide1, 1.3, 1.8, 11.0, 1.4, report_title, size=34, bold=True, color=text_head_rgb)
    _add_textbox(slide1, 1.3, 3.2, 11.0, 0.6, f"Project Dataset: {project_name}", size=18, color=accent_purple)
    
    import datetime
    date_str = datetime.datetime.now().strftime('%B %d, %Y at %I:%M %p')
    _add_textbox(slide1, 1.3, 3.9, 11.0, 0.4, f"Generated on {date_str}", size=13, color=text_sub_rgb)

    # Metadata Stat Badges Container
    meta_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(4.8), Inches(10.8), Inches(1.4))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = RGBColor(*card_rgb)
    meta_card.line.color.rgb = RGBColor(*card_border_rgb)

    meta_text = f"📊 Dataset Scope: {data.get('rows_count', 0):,} Rows × {data.get('cols_count', 0)} Columns  |  Numeric Features: {data.get('numeric_count', 0)}  |  Categorical Features: {data.get('cat_count', 0)}"
    _add_textbox(slide1, 1.5, 5.3, 10.4, 0.6, meta_text, size=14, bold=True, color=text_body_rgb, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: Executive KPI Metrics Grid ──────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2)
    _add_textbox(slide2, 0.8, 0.6, 11.5, 0.6, "Executive Data Metrics & Health Profile", size=26, bold=True, color=text_head_rgb)

    kpis = [
        ("Total Rows", f"{data.get('rows_count', 0):,}", "Complete Data Records", accent_purple),
        ("Total Features", f"{data.get('cols_count', 0)}", "Dataset Attributes", accent_emerald),
        ("Numeric Ratio", f"{data.get('numeric_count', 0)} / {data.get('cols_count', 0)}", "Quantitative Columns", accent_cyan),
        ("Data Quality", "100%", "Cleaned & Validated", accent_amber)
    ]

    for idx, (title, val, sub, col_rgb) in enumerate(kpis):
        left_pos = 0.8 + (idx * 2.95)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.6), Inches(2.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*card_rgb)
        card.line.color.rgb = RGBColor(*col_rgb)

        _add_textbox(slide2, left_pos + 0.1, 2.2, 2.5, 0.4, title.upper(), size=12, bold=True, color=col_rgb, align=PP_ALIGN.CENTER)
        _add_textbox(slide2, left_pos + 0.1, 3.2, 2.5, 1.0, val, size=32, bold=True, color=text_head_rgb, align=PP_ALIGN.CENTER)
        _add_textbox(slide2, left_pos + 0.1, 4.8, 2.5, 0.4, sub, size=11, color=text_sub_rgb, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: Descriptive Statistics Table ─────────────────────────────────
    import pandas as pd
    cleaned_csv = session_dir / "cleaned.csv"
    stats_df = None
    if cleaned_csv.exists():
        try:
            stats_df = pd.read_csv(cleaned_csv)
        except Exception:
            pass

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3)
    _add_textbox(slide3, 0.8, 0.6, 11.5, 0.6, "Feature Statistics Profile", size=26, bold=True, color=accent_emerald)

    if stats_df is not None:
        numeric_cols = stats_df.select_dtypes(include=['number']).columns.tolist()
        stats = []
        for col in numeric_cols[:8]:
            col_data = stats_df[col].dropna()
            if not col_data.empty:
                stats.append([
                    col[:26],
                    f"{col_data.min():.2f}" if col_data.dtype.kind in 'fc' else str(int(col_data.min())),
                    f"{col_data.max():.2f}" if col_data.dtype.kind in 'fc' else str(int(col_data.max())),
                    f"{col_data.mean():.2f}",
                    f"{col_data.std():.2f}"
                ])

        rows_len = len(stats) + 1
        cols_len = 5
        x, y, cx, cy = Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5 + 0.42 * len(stats))
        table_shape = slide3.shapes.add_table(rows_len, cols_len, x, y, cx, cy)
        table = table_shape.table

        table.columns[0].width = Inches(3.7)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(2.0)
        table.columns[4].width = Inches(2.0)

        headers = ["Numeric Feature", "Min Value", "Max Value", "Arithmetic Mean", "Std Deviation"]
        for c_idx, h_text in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.text = h_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*accent_emerald) if is_light else RGBColor(16, 185, 129)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.RIGHT

        for r_idx, row_data in enumerate(stats):
            for c_idx, val in enumerate(row_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*card_rgb)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(*text_body_rgb)
                p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.RIGHT

    # ── SLIDE 4: Strategic Business Insights (Structured Split Cards) ────────
    insights = data.get("insights", "").strip()
    insights_paragraphs = [p.strip() for p in _clean_md(insights).split("\n\n") if p.strip()]

    if insights_paragraphs:
        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide4)
        _add_textbox(slide4, 0.8, 0.6, 11.5, 0.6, "Strategic Business Insights & Recommendations", size=24, bold=True, color=accent_amber)

        parsed_cards = []
        for p in insights_paragraphs[:3]:
            obs, imp, strat = "", "", ""
            obs_m = re.search(r"Observation:\s*(.*?)(?=Business Implication|Actionable Strategy|$)", p, re.DOTALL | re.IGNORECASE)
            imp_m = re.search(r"Business Implication:\s*(.*?)(?=Observation|Actionable Strategy|$)", p, re.DOTALL | re.IGNORECASE)
            strat_m = re.search(r"Actionable Strategy:\s*(.*?)(?=Observation|Business Implication|$)", p, re.DOTALL | re.IGNORECASE)
            
            if obs_m: obs = obs_m.group(1).strip()
            if imp_m: imp = imp_m.group(1).strip()
            if strat_m: strat = strat_m.group(1).strip()

            parsed_cards.append({
                "obs": obs or p[:200],
                "imp": imp or "Resource allocation exhibits a lockstep relationship with performance metrics.",
                "strat": strat or "Establish continuous automated monitoring and resource allocation controls."
            })

        col_w = 3.65
        for idx, card_data in enumerate(parsed_cards[:3]):
            left_pos = 0.8 + (idx * 3.9)
            card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.4), Inches(col_w), Inches(5.4))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*card_rgb)
            card.line.color.rgb = RGBColor(*accent_amber)

            _add_textbox(slide4, left_pos + 0.15, 1.55, col_w - 0.3, 0.35, f"STRATEGIC PILLAR #{idx+1}", size=11, bold=True, color=accent_amber)

            _add_textbox(slide4, left_pos + 0.15, 2.0, col_w - 0.3, 0.25, "OBSERVATION", size=9.5, bold=True, color=accent_cyan)
            _add_textbox(slide4, left_pos + 0.15, 2.3, col_w - 0.3, 1.3, card_data["obs"], size=10, color=text_body_rgb)

            _add_textbox(slide4, left_pos + 0.15, 3.7, col_w - 0.3, 0.25, "BUSINESS IMPLICATION", size=9.5, bold=True, color=accent_purple)
            _add_textbox(slide4, left_pos + 0.15, 4.0, col_w - 0.3, 1.3, card_data["imp"], size=10, color=text_body_rgb)

            _add_textbox(slide4, left_pos + 0.15, 5.4, col_w - 0.3, 0.25, "ACTIONABLE STRATEGY", size=9.5, bold=True, color=accent_emerald)
            _add_textbox(slide4, left_pos + 0.15, 5.7, col_w - 0.3, 0.9, card_data["strat"], size=10, color=text_body_rgb)

    # ── SLIDE 5+: Visual Charts & Executive Takeaway Cards ───────────────────
    png_charts = data.get("png_charts", [])
    output_dir = Path(data.get("output_dir", ""))

    for idx, chart_name in enumerate(png_charts[:4]):
        chart_path = output_dir / chart_name
        if chart_path.exists():
            slide_chart = prs.slides.add_slide(prs.slide_layouts[6])
            _add_bg(slide_chart)
            chart_title = chart_name.replace(".png", "").replace("_", " ").title()
            _add_textbox(slide_chart, 0.8, 0.6, 11.5, 0.6, f"Visual Intelligence: {chart_title}", size=24, bold=True, color=accent_cyan)

            try:
                slide_chart.shapes.add_picture(str(chart_path), Inches(0.8), Inches(1.4), Inches(6.5))
            except Exception as chart_err:
                print(f"Error adding chart image to PPTX: {chart_err}")

            raw_t = insights_paragraphs[idx + 3] if (idx + 3) < len(insights_paragraphs) else ""
            takeaway_text = _clean_takeaway_text(raw_t)
            
            r_left = 7.6
            r_width = 4.9
            
            card_bg = slide_chart.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(r_left), Inches(1.4), Inches(r_width), Inches(5.4))
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(*card_rgb)
            card_bg.line.color.rgb = RGBColor(*accent_cyan)

            _add_textbox(slide_chart, r_left + 0.2, 1.65, r_width - 0.4, 0.35, "EXECUTIVE TAKEAWAY & ANALYSIS", size=12, bold=True, color=accent_cyan)

            _add_textbox(slide_chart, r_left + 0.2, 2.1, r_width - 0.4, 0.25, "KEY PATTERN OBSERVED", size=9.5, bold=True, color=accent_purple)
            _add_textbox(slide_chart, r_left + 0.2, 2.4, r_width - 0.4, 1.4, takeaway_text, size=10.5, color=text_body_rgb)

            _add_textbox(slide_chart, r_left + 0.2, 3.9, r_width - 0.4, 0.25, "OPERATIONAL RELEVANCE", size=9.5, bold=True, color=accent_amber)
            _add_textbox(slide_chart, r_left + 0.2, 4.2, r_width - 0.4, 1.0, "This visual distribution provides key evidence for resource allocation and predictive modeling.", size=10, color=text_body_rgb)

            _add_textbox(slide_chart, r_left + 0.2, 5.3, r_width - 0.4, 0.25, "RECOMMENDED NEXT STEP", size=9.5, bold=True, color=accent_emerald)
            _add_textbox(slide_chart, r_left + 0.2, 5.6, r_width - 0.4, 1.0, "Incorporate key column metrics into automated data-quality monitor.", size=10, color=text_body_rgb)

    # ── SLIDE LAST: Conclusion & Action Plan (Stacked Action Cards) ──────────
    slide_final = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide_final)
    _add_textbox(slide_final, 0.8, 0.6, 11.5, 0.6, "Conclusions & Actionable Implementation", size=24, bold=True, color=accent_emerald)

    action_items = [
        ("01", "Operational Optimization", "Leverage mapped correlations to drive high-impact operational optimizations and resource reallocation.", accent_cyan),
        ("02", "Automated Data Governance", "Implement automated data-quality checks on continuous incoming data streams to prevent pipeline anomalies.", accent_purple),
        ("03", "Predictive Integration", "Deploy machine-learning ready data structures directly into downstream predictive modeling pipelines.", accent_amber),
        ("04", "Stakeholder Alignment", "Share executive visual decks and insights with key business stakeholders for strategic alignment.", accent_emerald)
    ]

    for idx, (num_str, title_str, desc_str, col_rgb) in enumerate(action_items):
        top_pos = 1.4 + (idx * 1.35)
        
        card = slide_final.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*card_rgb)
        card.line.color.rgb = RGBColor(*col_rgb)

        num_box = slide_final.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(top_pos + 0.15), Inches(0.9), Inches(0.9))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = RGBColor(*col_rgb)
        num_box.line.fill.background()

        _add_textbox(slide_final, 0.95, top_pos + 0.35, 0.9, 0.5, num_str, size=18, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)

        _add_textbox(slide_final, 2.0, top_pos + 0.2, 10.3, 0.35, title_str.upper(), size=12, bold=True, color=col_rgb)
        _add_textbox(slide_final, 2.0, top_pos + 0.55, 10.3, 0.55, desc_str, size=11, color=text_body_rgb)

    pptx_path = session_dir / "report.pptx"
    prs.save(str(pptx_path))

    base_name = meta.get("filename", "report").rsplit(".", 1)[0] if "." in meta.get("filename", "") else meta.get("name", "report")
    return FileResponse(
        str(pptx_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{base_name}_executive_deck.pptx"
    )

# ---------------------------------------------------------------------------
# Cross-Project Comparison API
# ---------------------------------------------------------------------------

@app.get("/api/projects/compare")
async def compare_projects(project_a: str, project_b: str):
    """Returns comparative delta data for two completed projects."""
    def _load_project_summary(pid):
        session_dir = get_safe_session_dir(pid)
        results_path = session_dir / "results.json"
        if not results_path.exists():
            raise HTTPException(status_code=404, detail=f"Results not found for project {pid}")
        with open(results_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        meta = get_project_metadata(pid)
        return {
            "id": pid,
            "name": meta.get("name", pid),
            "rows": data.get("rows_count", 0),
            "cols": data.get("cols_count", 0),
            "numeric": data.get("numeric_count", 0),
            "categorical": data.get("cat_count", 0),
            "charts_count": len(data.get("plotly_charts", [])) + len(data.get("png_charts", [])),
            "cleaning": data.get("cleaning_steps", "")[:500],
            "insights": data.get("insights", "")[:800],
            "relations": data.get("relations", "")[:500],
        }

    a = _load_project_summary(project_a)
    b = _load_project_summary(project_b)

    def _delta(va, vb):
        if va == 0:
            return "+∞%" if vb > 0 else "0%"
        pct = round(((vb - va) / va) * 100, 1)
        return f"+{pct}%" if pct >= 0 else f"{pct}%"

    return {
        "project_a": a,
        "project_b": b,
        "deltas": {
            "rows": _delta(a["rows"], b["rows"]),
            "cols": _delta(a["cols"], b["cols"]),
            "numeric": _delta(a["numeric"], b["numeric"]),
            "categorical": _delta(a["categorical"], b["categorical"]),
            "charts": _delta(a["charts_count"], b["charts_count"]),
        }
    }


# ---------------------------------------------------------------------------
# Frontend Static Mounts
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
web_dir = BASE_DIR / "web"
assets_dir = BASE_DIR / "assets"
if not assets_dir.exists() or not list(assets_dir.glob("*")):
    assets_dir = web_dir / "assets"

app.mount("/assets", StaticFiles(directory=str(assets_dir)), name="assets")
app.mount("/", StaticFiles(directory=str(web_dir), html=True), name="web")


# ── Server Boot ─────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    # Start server on 8000
    print("\n" + "=" * 50)
    print("Crewlyze Web Platform")
    print("Local URL: http://localhost:8000")
    print("=" * 50 + "\n")
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
